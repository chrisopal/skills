#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from ppt_renderer import render_ppt_from_specs


SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Assemble slide assets into a PPTX file.")
    parser.add_argument("--images", nargs="*", help="Ordered list of slide image paths")
    parser.add_argument("--slide-specs", help="Path to slide_specs.json for direct editable rendering")
    parser.add_argument("--master-style", help="Path to master_style.json when using --slide-specs")
    parser.add_argument("--output", required=True, help="Output .pptx path")
    return parser


def validate_images(image_paths: list[Path]) -> None:
    for path in image_paths:
        if not path.exists():
            raise FileNotFoundError(f"Image not found: {path}")
        with Image.open(path) as image:
            image.verify()


def assemble_image_pptx(image_paths: list[Path], output_path: Path) -> Path:
    validate_images(image_paths)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    blank_layout = prs.slide_layouts[6]

    for image_path in image_paths:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(image_path),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def assemble_pptx(image_paths: list[Path], output_path: Path) -> Path:
    return assemble_image_pptx(image_paths, output_path)


def assemble_from_specs(slide_specs_path: Path, master_style_path: Path, output_path: Path, *, js_generator=None) -> Path:
    slide_specs = json.loads(slide_specs_path.read_text(encoding="utf-8"))
    master_style = json.loads(master_style_path.read_text(encoding="utf-8"))
    return render_ppt_from_specs(slide_specs, master_style, output_path, js_generator=js_generator)


def main() -> int:
    parser = build_parser()
    args = parser.parse_args()

    output_path = Path(args.output).expanduser().resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)

    if args.slide_specs:
        if not args.master_style:
            parser.error("--master-style is required when using --slide-specs")
        slide_specs_path = Path(args.slide_specs).expanduser().resolve()
        master_style_path = Path(args.master_style).expanduser().resolve()
        assemble_from_specs(slide_specs_path, master_style_path, output_path)
        print(f"[OK] Wrote PPTX: {output_path}")
        return 0

    if not args.images:
        parser.error("Provide --images or --slide-specs")

    image_paths = [Path(p).expanduser().resolve() for p in args.images]
    assemble_image_pptx(image_paths, output_path)
    print(f"[OK] Wrote PPTX: {output_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
