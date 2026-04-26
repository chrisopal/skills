"""Deterministic PPTX renderers for the 12 shipped patterns.

Each renderer takes (slide, slots, regions, theme) and draws the page using
python-pptx native shapes. No LLM, no Node — purely deterministic + offline.
Slot data flows straight into shapes so what the user typed is what they get.

Theme is a flat color dict pulled from master_style.color_strategy with
sensible fallbacks. Regions come straight from the pattern definition (or
from a sanitized fallback).
"""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

DEFAULT_THEME = {
    "primary": "#A8D86B",
    "secondary": "#0F95B6",
    "background": "#FFFFFF",
    "panel": "#F5F7FA",
    "text_primary": "#1E1E1E",
    "text_secondary": "#6B7280",
    "divider": "#E5E7EB",
}


def theme_from_master_style(master_style: dict | None) -> dict:
    if not master_style:
        return dict(DEFAULT_THEME)
    cs = master_style.get("color_strategy") or {}
    return {
        "primary": cs.get("primary_green") or cs.get("primary") or cs.get("primary_blue") or cs.get("accent_green") or DEFAULT_THEME["primary"],
        "secondary": cs.get("secondary_teal") or cs.get("secondary") or cs.get("accent_amber") or DEFAULT_THEME["secondary"],
        "background": cs.get("background") or DEFAULT_THEME["background"],
        "panel": cs.get("section_background") or cs.get("panel_background") or DEFAULT_THEME["panel"],
        "text_primary": cs.get("text_primary") or DEFAULT_THEME["text_primary"],
        "text_secondary": cs.get("text_secondary") or DEFAULT_THEME["text_secondary"],
        "divider": cs.get("divider") or DEFAULT_THEME["divider"],
    }


def _expand_hex(value: str) -> str:
    if len(value) == 4 and value.startswith("#"):
        return "#" + "".join(ch * 2 for ch in value[1:])
    return value


def _rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(_expand_hex(hex_str).lstrip("#").upper())


def _add_text(
    slide,
    text: str,
    x: float, y: float, w: float, h: float,
    *,
    size: int = 14,
    bold: bool = False,
    color: str = "#1E1E1E",
    align: str | None = None,
):
    if not text:
        return None
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return box


def _add_card(
    slide,
    x: float, y: float, w: float, h: float,
    *,
    fill: str = "#FFFFFF",
    stroke: str | None = None,
    stroke_pt: float = 1.0,
):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    if stroke:
        shape.line.color.rgb = _rgb(stroke)
        shape.line.width = Pt(stroke_pt)
    else:
        shape.line.fill.background()
    return shape


def _add_band(slide, x, y, w, h, *, fill: str):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.fill.background()
    return shape


def _add_line(slide, x1, y1, x2, y2, *, color: str, weight_pt: float = 1.0):
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE

    connector = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = _rgb(color)
    connector.line.width = Pt(weight_pt)
    return connector


def _safe(slots: dict, key: str, default: str = "") -> str:
    value = slots.get(key, default) if isinstance(slots, dict) else default
    if value is None:
        return default
    return str(value)


# ---------------------------------------------------------------------------
# Per-pattern renderers
# ---------------------------------------------------------------------------


def render_cover(slide, slots, regions, theme):
    _add_text(slide, _safe(slots, "title"), 0.6, 2.2, 12.0, 1.6,
              size=44, bold=True, color=theme["text_primary"])
    _add_text(slide, _safe(slots, "subtitle"), 0.6, 3.9, 12.0, 0.8,
              size=22, color=theme["text_secondary"])
    _add_line(slide, 0.6, 3.7, 6.5, 3.7, color=theme["primary"], weight_pt=2.0)
    _add_text(slide, _safe(slots, "org_block"), 0.6, 6.5, 12.0, 0.5,
              size=14, color=theme["text_secondary"])


def render_section_divider(slide, slots, regions, theme):
    section_no = _safe(slots, "section_no", "01")
    section_title = _safe(slots, "section_title")
    agenda = _safe(slots, "agenda_anchor")
    _add_text(slide, section_no, 0.6, 1.4, 6.0, 2.5,
              size=160, bold=True, color=theme["primary"])
    _add_line(slide, 0.6, 4.0, 12.7, 4.0, color=theme["secondary"], weight_pt=2.0)
    _add_text(slide, section_title, 0.6, 4.3, 12.0, 0.9,
              size=36, bold=True, color=theme["text_primary"])
    if agenda:
        _add_text(slide, agenda, 0.6, 5.3, 12.0, 0.6,
                  size=18, color=theme["text_secondary"])


def render_conclusion_top_modules(slide, slots, regions, theme):
    _add_band(slide, 0.4, 0.4, 12.5, 1.3, fill=theme["panel"])
    _add_text(slide, _safe(slots, "headline"), 0.7, 0.55, 12.0, 1.0,
              size=24, bold=True, color=theme["text_primary"])
    cols = []
    for idx in (1, 2, 3, 4):
        title = _safe(slots, f"module_{idx}_title")
        body = _safe(slots, f"module_{idx}_body")
        if not title and not body:
            continue
        cols.append((idx, title, body))
    if not cols:
        return
    n = len(cols)
    total_w = 12.5
    gap = 0.2
    each_w = (total_w - gap * (n - 1)) / n
    for i, (idx, title, body) in enumerate(cols):
        x = 0.4 + i * (each_w + gap)
        y = 2.0
        h = 5.0
        _add_card(slide, x, y, each_w, h, stroke=theme["primary"])
        _add_text(slide, title, x + 0.2, y + 0.3, each_w - 0.4, 0.7,
                  size=20, bold=True, color=theme["text_primary"])
        _add_text(slide, body, x + 0.2, y + 1.1, each_w - 0.4, h - 1.4,
                  size=14, color=theme["text_primary"])


def render_two_column_compare(slide, slots, regions, theme):
    _add_text(slide, "对照分析", 0.5, 0.4, 12.3, 0.9, size=24, bold=True, color=theme["text_primary"])
    # Left column
    _add_card(slide, 0.5, 1.5, 6.0, 4.5, stroke=theme["primary"])
    _add_text(slide, _safe(slots, "left_title"), 0.7, 1.7, 5.6, 0.7,
              size=22, bold=True, color=theme["text_primary"])
    for i, key in enumerate(("left_point_1", "left_point_2", "left_point_3"), start=0):
        text = _safe(slots, key)
        if not text:
            continue
        _add_text(slide, f"• {text}", 0.7, 2.6 + i * 0.9, 5.6, 0.85,
                  size=16, color=theme["text_primary"])
    # Right column
    _add_card(slide, 6.83, 1.5, 6.0, 4.5, stroke=theme["secondary"])
    _add_text(slide, _safe(slots, "right_title"), 7.05, 1.7, 5.6, 0.7,
              size=22, bold=True, color=theme["text_primary"])
    for i, key in enumerate(("right_point_1", "right_point_2", "right_point_3"), start=0):
        text = _safe(slots, key)
        if not text:
            continue
        _add_text(slide, f"• {text}", 7.05, 2.6 + i * 0.9, 5.6, 0.85,
                  size=16, color=theme["text_primary"])
    verdict = _safe(slots, "verdict")
    if verdict:
        _add_band(slide, 0.5, 6.3, 12.3, 0.85, fill=theme["panel"])
        _add_text(slide, verdict, 0.7, 6.4, 12.0, 0.7,
                  size=18, bold=True, color=theme["secondary"])


def render_four_card_matrix(slide, slots, regions, theme):
    _add_text(slide, "核心指标矩阵", 0.5, 0.4, 12.3, 0.9, size=24, bold=True, color=theme["text_primary"])
    cells = [
        (1, 0.5, 1.6),
        (2, 6.83, 1.6),
        (3, 0.5, 4.4),
        (4, 6.83, 4.4),
    ]
    for idx, x, y in cells:
        label = _safe(slots, f"cell_{idx}_label")
        value = _safe(slots, f"cell_{idx}_value")
        desc = _safe(slots, f"cell_{idx}_desc")
        _add_card(slide, x, y, 6.0, 2.6, stroke=theme["primary"])
        _add_text(slide, label, x + 0.25, y + 0.2, 5.5, 0.4,
                  size=14, color=theme["text_secondary"])
        _add_text(slide, value, x + 0.25, y + 0.7, 5.5, 1.0,
                  size=44, bold=True, color=theme["primary"])
        _add_text(slide, desc, x + 0.25, y + 1.8, 5.5, 0.7,
                  size=13, color=theme["text_primary"])


def render_three_stage_path(slide, slots, regions, theme):
    _add_text(slide, "推进路线图", 0.5, 0.4, 12.3, 0.9, size=24, bold=True, color=theme["text_primary"])
    _add_line(slide, 1.2, 4.2, 12.1, 4.2, color=theme["secondary"], weight_pt=2.0)
    centers = [(2.4, 1), (6.65, 2), (10.9, 3)]
    for x, idx in centers:
        _add_card(slide, x - 0.4, 4.0, 0.8, 0.4, fill=theme["primary"])
        _add_text(slide, str(idx), x - 0.4, 4.0, 0.8, 0.4,
                  size=18, bold=True, color="#FFFFFF", align="center")
        phase = _safe(slots, f"stage_{idx}_phase")
        title = _safe(slots, f"stage_{idx}_title")
        deliverable = _safe(slots, f"stage_{idx}_deliverable")
        _add_card(slide, x - 1.6, 1.7, 3.2, 2.0, stroke=theme["primary"])
        _add_text(slide, phase, x - 1.4, 1.85, 2.8, 0.4, size=12, color=theme["text_secondary"])
        _add_text(slide, title, x - 1.4, 2.3, 2.8, 0.6, size=18, bold=True, color=theme["text_primary"])
        _add_text(slide, deliverable, x - 1.6, 4.7, 3.2, 1.5, size=13, color=theme["text_primary"])


def render_kpi_strip(slide, slots, regions, theme):
    _add_text(slide, "关键指标", 0.5, 0.4, 12.3, 0.9, size=24, bold=True, color=theme["text_primary"])
    cards_y = 1.7
    card_w = 3.0
    gap = 0.15
    start_x = 0.5
    for idx in (1, 2, 3, 4):
        value = _safe(slots, f"kpi_{idx}_value")
        if not value:
            continue
        x = start_x + (idx - 1) * (card_w + gap)
        _add_card(slide, x, cards_y, card_w, 4.0, stroke=theme["primary"])
        label = _safe(slots, f"kpi_{idx}_label")
        unit = _safe(slots, f"kpi_{idx}_unit")
        trend = _safe(slots, f"kpi_{idx}_trend")
        _add_text(slide, label, x + 0.2, cards_y + 0.25, card_w - 0.4, 0.5,
                  size=14, color=theme["text_secondary"])
        _add_text(slide, value, x + 0.2, cards_y + 0.95, card_w - 0.4, 1.4,
                  size=52, bold=True, color=theme["primary"])
        _add_text(slide, unit, x + 0.2, cards_y + 2.5, card_w - 0.4, 0.4,
                  size=14, color=theme["text_secondary"])
        _add_text(slide, trend, x + 0.2, cards_y + 3.05, card_w - 0.4, 0.6,
                  size=13, color=theme["secondary"])


def render_architecture_layers(slide, slots, regions, theme):
    _add_text(slide, "架构分层", 0.5, 0.4, 12.3, 0.9, size=24, bold=True, color=theme["text_primary"])
    layers = []
    for idx in range(1, 6):
        name = _safe(slots, f"layer_{idx}_name")
        components = _safe(slots, f"layer_{idx}_components")
        if not name and not components:
            continue
        layers.append((idx, name, components))
    if not layers:
        return
    band_h = (5.5 / max(len(layers), 1)) - 0.1
    for i, (idx, name, components) in enumerate(layers):
        y = 1.6 + i * (band_h + 0.1)
        _add_band(slide, 0.5, y, 12.3, band_h, fill=theme["panel"])
        _add_text(slide, name, 0.7, y + 0.1, 4.0, band_h - 0.1,
                  size=18, bold=True, color=theme["secondary"])
        _add_text(slide, components, 4.8, y + 0.1, 7.9, band_h - 0.1,
                  size=14, color=theme["text_primary"])


def render_before_after(slide, slots, regions, theme):
    _add_text(slide, "前后对比", 0.5, 0.4, 12.3, 0.9, size=24, bold=True, color=theme["text_primary"])
    # Before
    _add_card(slide, 0.5, 1.5, 6.0, 4.0, fill=theme["panel"], stroke=theme["divider"])
    _add_text(slide, "Before", 0.7, 1.65, 5.6, 0.4, size=14, color=theme["text_secondary"])
    _add_text(slide, _safe(slots, "before_title"), 0.7, 2.05, 5.6, 0.7,
              size=22, bold=True, color=theme["text_primary"])
    for i, key in enumerate(("before_point_1", "before_point_2", "before_point_3")):
        text = _safe(slots, key)
        if text:
            _add_text(slide, f"• {text}", 0.7, 2.85 + i * 0.6, 5.6, 0.55,
                      size=14, color=theme["text_primary"])
    # After
    _add_card(slide, 6.83, 1.5, 6.0, 4.0, fill=theme["background"], stroke=theme["primary"])
    _add_text(slide, "After", 7.05, 1.65, 5.6, 0.4, size=14, color=theme["primary"])
    _add_text(slide, _safe(slots, "after_title"), 7.05, 2.05, 5.6, 0.7,
              size=22, bold=True, color=theme["text_primary"])
    for i, key in enumerate(("after_point_1", "after_point_2", "after_point_3")):
        text = _safe(slots, key)
        if text:
            _add_text(slide, f"• {text}", 7.05, 2.85 + i * 0.6, 5.6, 0.55,
                      size=14, color=theme["text_primary"])
    delta = _safe(slots, "delta_callout")
    if delta:
        _add_band(slide, 0.5, 5.85, 12.3, 0.9, fill=theme["panel"])
        _add_text(slide, delta, 0.7, 5.95, 12.0, 0.7,
                  size=18, bold=True, color=theme["secondary"])


def render_evidence_grid(slide, slots, regions, theme):
    _add_text(slide, "案例与证据", 0.5, 0.4, 12.3, 0.9, size=24, bold=True, color=theme["text_primary"])
    positions = [
        (0.5, 1.5), (4.78, 1.5), (9.06, 1.5),
        (0.5, 4.5), (4.78, 4.5), (9.06, 4.5),
    ]
    for idx, (x, y) in enumerate(positions, start=1):
        headline = _safe(slots, f"case_{idx}_headline")
        metric = _safe(slots, f"case_{idx}_metric")
        if not headline and not metric:
            continue
        _add_card(slide, x, y, 4.0, 2.7, stroke=theme["primary"])
        _add_text(slide, headline, x + 0.2, y + 0.3, 3.6, 1.2,
                  size=16, bold=True, color=theme["text_primary"])
        _add_text(slide, metric, x + 0.2, y + 1.6, 3.6, 1.0,
                  size=28, bold=True, color=theme["primary"])


def render_summary_takeaways(slide, slots, regions, theme):
    _add_text(slide, "下一步行动", 0.5, 0.4, 12.3, 0.9, size=24, bold=True, color=theme["text_primary"])
    takeaways = []
    for idx in range(1, 6):
        text = _safe(slots, f"takeaway_{idx}")
        if text:
            takeaways.append((idx, text))
    band_h = 0.85
    start_y = 1.7
    for i, (idx, text) in enumerate(takeaways):
        y = start_y + i * (band_h + 0.1)
        _add_card(slide, 0.5, y, 12.3, band_h, stroke=theme["primary"])
        # Number badge
        _add_card(slide, 0.7, y + 0.18, 0.5, 0.5, fill=theme["primary"])
        _add_text(slide, str(idx), 0.7, y + 0.18, 0.5, 0.5,
                  size=14, bold=True, color="#FFFFFF", align="center")
        _add_text(slide, text, 1.4, y + 0.18, 10.6, band_h - 0.2,
                  size=18, color=theme["text_primary"])
    next_step = _safe(slots, "next_step")
    if next_step:
        next_y = start_y + len(takeaways) * (band_h + 0.1) + 0.2
        if next_y < 6.6:
            _add_band(slide, 0.5, next_y, 12.3, 0.7, fill=theme["panel"])
            _add_text(slide, next_step, 0.7, next_y + 0.05, 12.0, 0.6,
                      size=18, color=theme["secondary"])


def render_freeform(slide, slots, regions, theme):
    headline = _safe(slots, "headline")
    note = _safe(slots, "note")
    _add_text(slide, headline, 0.5, 0.6, 12.3, 1.2, size=32, bold=True, color=theme["text_primary"])
    _add_text(slide, note, 0.5, 2.0, 12.3, 4.5, size=16, color=theme["text_secondary"])
    _add_text(slide, "[ freeform — supply layout_design + grid in page intent ]",
              0.5, 6.7, 12.3, 0.4, size=11, color=theme["divider"])


PATTERN_RENDERERS = {
    "cover": render_cover,
    "section_divider": render_section_divider,
    "conclusion_top_modules": render_conclusion_top_modules,
    "two_column_compare": render_two_column_compare,
    "four_card_matrix": render_four_card_matrix,
    "three_stage_path": render_three_stage_path,
    "kpi_strip": render_kpi_strip,
    "architecture_layers": render_architecture_layers,
    "before_after": render_before_after,
    "evidence_grid": render_evidence_grid,
    "summary_takeaways": render_summary_takeaways,
    "freeform": render_freeform,
}


def render_slide(slide, pattern_id: str, slots: dict, regions: dict | None, theme: dict) -> bool:
    """Dispatch to the matching renderer; return True iff a renderer was found."""

    func = PATTERN_RENDERERS.get(pattern_id)
    if func is None:
        return False
    func(slide, slots or {}, regions or {}, theme)
    return True
