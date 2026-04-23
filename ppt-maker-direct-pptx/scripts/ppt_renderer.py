#!/usr/bin/env python3
from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5


@dataclass
class Theme:
    background: str
    surface: str
    primary: str
    accent: str
    neutral: str
    text: str
    text_muted: str
    divider: str
    title_font: str
    body_font: str


def rgb(value: str | None, fallback: str) -> RGBColor:
    raw = (value or fallback).replace("#", "").strip()
    if len(raw) == 3:
        raw = "".join(ch * 2 for ch in raw)
    return RGBColor.from_string(raw.upper())


def to_theme(master_style: dict[str, Any]) -> Theme:
    colors = master_style.get("color_strategy", {})
    typography = master_style.get("typography", {})
    return Theme(
        background=colors.get("background", "#FFFFFF"),
        surface=colors.get("section_background", "#F5F7FA"),
        primary=colors.get("secondary_teal") or colors.get("primary_blue") or "#0F95B6",
        accent=colors.get("primary_green") or colors.get("accent_green") or "#A8D86B",
        neutral=colors.get("neutral_gray", "#D9D9D9"),
        text=colors.get("text_primary", "#1E1E1E"),
        text_muted=colors.get("text_secondary", "#6B7280"),
        divider=colors.get("divider", "#E5E7EB"),
        title_font=typography.get("title_font", "Microsoft YaHei"),
        body_font=typography.get("body_font", "Microsoft YaHei"),
    )


def split_summary(text: str, limit: int = 36) -> list[str]:
    text = text.strip()
    if not text:
        return []
    if len(text) <= limit:
        return [text]
    words = text.replace("；", "，").split("，")
    lines: list[str] = []
    current = ""
    for word in words:
        word = word.strip()
        if not word:
            continue
        candidate = f"{current}｜{word}" if current else word
        if len(candidate) > limit and current:
            lines.append(current)
            current = word
        else:
            current = candidate
    if current:
        lines.append(current)
    return lines[:3]


def normalize_blocks(key_blocks: list[Any]) -> list[dict[str, Any]]:
    normalized: list[dict[str, Any]] = []
    for idx, block in enumerate(key_blocks, start=1):
        if isinstance(block, dict):
            normalized.append(
                {
                    "title": str(block.get("title") or f"模块 {idx}").strip(),
                    "items": [str(item).strip() for item in block.get("items", []) if str(item).strip()],
                    "summary": str(block.get("summary", "")).strip(),
                }
            )
        else:
            text = str(block).strip()
            normalized.append({"title": text or f"模块 {idx}", "items": [], "summary": ""})
    return normalized


def infer_layout_family(slide: dict[str, Any], template_name: str) -> str:
    text = " ".join(
        str(slide.get(key, "") or "")
        for key in ("title", "subtitle", "page_goal", "layout_type", "visual_focus", "detail_notes")
    ).lower()
    block_count = len(normalize_blocks(slide.get("key_blocks", [])))

    if any(word in text for word in ["架构", "蓝图", "底座", "层", "平台", "architecture", "blueprint"]):
        return "architecture_layers"
    if any(word in text for word in ["路线", "里程碑", "阶段", "实施", "roadmap", "milestone", "计划"]):
        return "roadmap"
    if block_count >= 4:
        return "four_card_grid"
    if block_count == 3:
        return "three_card_grid"
    if block_count == 2:
        return "two_column"
    if any(word in text for word in ["总结", "收益", "价值", "结论", "建议", "summary", "value"]):
        return "summary"
    if "会议" in template_name or "内部" in template_name:
        return "two_column"
    return "focus_with_sidebar"


def build_slide_specs(job: dict[str, Any], master_style: dict[str, Any], slide_prompts: dict[str, Any]) -> dict[str, Any]:
    slides = []
    template_name = job.get("template_name") or job.get("style") or ""
    for slide in slide_prompts.get("slides", []):
        blocks = normalize_blocks(slide.get("key_blocks", []))
        slides.append(
            {
                "page_no": slide.get("page_no"),
                "title": slide.get("title", ""),
                "subtitle": slide.get("subtitle", ""),
                "page_goal": slide.get("page_goal", ""),
                "visual_focus": slide.get("visual_focus", ""),
                "detail_notes": slide.get("detail_notes", ""),
                "layout_family": infer_layout_family(slide, template_name),
                "blocks": blocks,
                "compiled_prompt": slide.get("compiled_prompt", ""),
            }
        )
    return {"slides": slides}


def add_textbox(
    slide,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    font_face: str,
    font_size: float,
    color: str,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_face
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color, "1E1E1E")
    return box


def add_bullets(
    slide,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    items: list[str],
    theme: Theme,
    font_size: float = 12,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for index, item in enumerate(items[:5]):
        para = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.bullet = True
        run = para.add_run()
        run.text = item
        run.font.name = theme.body_font
        run.font.size = Pt(font_size)
        run.font.color.rgb = rgb(theme.text, "1E1E1E")
    return box


def add_card(slide, *, x: float, y: float, w: float, h: float, theme: Theme, accent: str, title: str, items: list[str], summary: str = ""):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(theme.surface, "F5F7FA")
    shape.line.color.rgb = rgb(theme.divider, "E5E7EB")
    shape.line.width = Pt(1.0)

    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(0.08),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = rgb(accent, "0F95B6")
    band.line.fill.background()

    add_textbox(
        slide,
        x=x + 0.12,
        y=y + 0.14,
        w=w - 0.24,
        h=0.28,
        text=title,
        font_face=theme.body_font,
        font_size=14,
        color=accent,
        bold=True,
    )

    body_items = items[:4]
    if summary and not body_items:
        body_items = split_summary(summary, 28)
    add_bullets(slide, x=x + 0.12, y=y + 0.46, w=w - 0.24, h=h - 0.58, items=body_items, theme=theme, font_size=11.5)


def add_header(slide, spec: dict[str, Any], theme: Theme):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM,
        Inches(0.45),
        Inches(0.28),
        Inches(0.62),
        Inches(0.24),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(theme.primary, "0F95B6")
    shape.line.fill.background()
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM,
        Inches(0.98),
        Inches(0.28),
        Inches(0.40),
        Inches(0.24),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(theme.accent, "A8D86B")
    shape.line.fill.background()

    add_textbox(
        slide,
        x=0.72,
        y=0.34,
        w=8.5,
        h=0.36,
        text=spec["title"],
        font_face=theme.title_font,
        font_size=24,
        color=theme.text,
        bold=True,
    )
    if spec.get("subtitle"):
        add_textbox(
            slide,
            x=0.74,
            y=0.76,
            w=8.8,
            h=0.24,
            text=spec["subtitle"],
            font_face=theme.body_font,
            font_size=12,
            color=theme.text_muted,
        )

    chip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(10.65),
        Inches(0.40),
        Inches(1.9),
        Inches(0.34),
    )
    chip.fill.solid()
    chip.fill.fore_color.rgb = rgb(theme.surface, "F5F7FA")
    chip.line.color.rgb = rgb(theme.primary, "0F95B6")
    chip.line.width = Pt(1.1)
    add_textbox(
        slide,
        x=10.78,
        y=0.47,
        w=1.6,
        h=0.18,
        text=spec["layout_family"],
        font_face=theme.body_font,
        font_size=10,
        color=theme.primary,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def render_architecture_layers(slide, spec: dict[str, Any], theme: Theme):
    blocks = spec["blocks"][:5]
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.56),
        Inches(1.32),
        Inches(12.18),
        Inches(5.76),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = rgb("FBFCFD", "FBFCFD")
    panel.line.color.rgb = rgb(theme.divider, "E5E7EB")
    panel.line.width = Pt(1.0)

    label_x = 0.76
    row_x = 2.05
    row_w = 10.25
    row_h = 0.78
    base_y = 1.55

    for idx, block in enumerate(blocks):
        y = base_y + idx * 0.92
        label = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(label_x),
            Inches(y),
            Inches(1.18),
            Inches(0.56),
        )
        label.fill.solid()
        label.fill.fore_color.rgb = rgb(theme.primary, "0F95B6")
        label.line.fill.background()
        add_textbox(
            slide,
            x=label_x + 0.12,
            y=y + 0.12,
            w=0.94,
            h=0.16,
            text=block["title"],
            font_face=theme.body_font,
            font_size=11,
            color="FFFFFF",
            bold=True,
        )

        container = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(row_x),
            Inches(y - 0.05),
            Inches(row_w),
            Inches(row_h),
        )
        container.fill.solid()
        container.fill.fore_color.rgb = rgb(theme.surface, "F5F7FA")
        container.line.color.rgb = rgb(theme.divider, "E5E7EB")
        container.line.width = Pt(0.8)

        items = block["items"][:5] if block["items"] else split_summary(block["summary"] or block["title"], 14)
        cursor_x = row_x + 0.16
        widths = [1.8, 1.55, 1.55, 1.75, 1.95]
        for j, item in enumerate(items):
            width = widths[j] if j < len(widths) else 1.6
            accent = theme.primary if j < max(1, len(items) - 1) else theme.accent
            fill = theme.surface if accent == theme.accent else "FFFFFF"
            add_card(
                slide,
                x=cursor_x,
                y=y + 0.06,
                w=width,
                h=0.56,
                theme=theme,
                accent=accent,
                title=item,
                items=[],
                summary="",
            )
            cursor_x += width + 0.12

        if idx < len(blocks) - 1:
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(6.2),
                Inches(y + 0.73),
                Inches(6.2),
                Inches(y + 0.88),
            )
            connector.line.color.rgb = rgb(theme.primary, "0F95B6")
            connector.line.width = Pt(1.2)


def render_two_column(slide, spec: dict[str, Any], theme: Theme):
    blocks = spec["blocks"]
    left = blocks[0] if blocks else {"title": "模块一", "items": [], "summary": ""}
    right = blocks[1] if len(blocks) > 1 else {"title": "模块二", "items": [], "summary": ""}

    add_card(slide, x=0.72, y=1.62, w=5.72, h=4.75, theme=theme, accent=theme.primary, title=left["title"], items=left["items"], summary=left["summary"])
    add_card(slide, x=6.86, y=1.62, w=5.72, h=4.75, theme=theme, accent=theme.accent, title=right["title"], items=right["items"], summary=right["summary"])

    if spec.get("page_goal"):
        add_textbox(
            slide,
            x=0.72,
            y=6.56,
            w=11.86,
            h=0.24,
            text=f"页面目标：{spec['page_goal']}",
            font_face=theme.body_font,
            font_size=11,
            color=theme.text_muted,
        )


def render_card_grid(slide, spec: dict[str, Any], theme: Theme, columns: int):
    blocks = spec["blocks"][: max(3, columns * 2)]
    gap = 0.18
    x0 = 0.72
    y0 = 1.72
    total_w = 11.84
    card_w = (total_w - gap * (columns - 1)) / columns
    rows = 2 if len(blocks) > columns else 1
    card_h = 2.05 if rows == 2 else 3.9
    for idx, block in enumerate(blocks):
        row = idx // columns
        col = idx % columns
        accent = theme.accent if idx % 2 else theme.primary
        add_card(
            slide,
            x=x0 + col * (card_w + gap),
            y=y0 + row * (card_h + 0.24),
            w=card_w,
            h=card_h,
            theme=theme,
            accent=accent,
            title=block["title"],
            items=block["items"],
            summary=block["summary"],
        )


def render_summary(slide, spec: dict[str, Any], theme: Theme):
    statement = spec.get("page_goal") or spec.get("visual_focus") or "总结本页核心结论与收益。"
    banner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.72),
        Inches(1.62),
        Inches(11.84),
        Inches(0.74),
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = rgb(theme.primary, "0F95B6")
    banner.line.fill.background()
    add_textbox(
        slide,
        x=0.96,
        y=1.82,
        w=11.2,
        h=0.22,
        text=statement,
        font_face=theme.title_font,
        font_size=16,
        color="FFFFFF",
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    render_card_grid(slide, spec, theme, columns=3)


def render_roadmap(slide, spec: dict[str, Any], theme: Theme):
    blocks = spec["blocks"][:4]
    x0 = 0.92
    y = 2.6
    gap = 0.22
    card_w = 2.75
    card_h = 2.3
    for idx, block in enumerate(blocks):
        x = x0 + idx * (card_w + gap)
        add_card(slide, x=x, y=y, w=card_w, h=card_h, theme=theme, accent=theme.primary if idx != len(blocks) - 1 else theme.accent, title=block["title"], items=block["items"], summary=block["summary"])
        if idx < len(blocks) - 1:
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.CHEVRON,
                Inches(x + card_w),
                Inches(y + 1.02),
                Inches(x + card_w + gap - 0.04),
                Inches(y + 1.02),
            )
            connector.line.color.rgb = rgb(theme.primary, "0F95B6")
            connector.line.width = Pt(1.2)
    if spec.get("detail_notes"):
        add_textbox(
            slide,
            x=0.92,
            y=5.25,
            w=11.5,
            h=0.26,
            text=spec["detail_notes"],
            font_face=theme.body_font,
            font_size=11,
            color=theme.text_muted,
        )


def render_focus_with_sidebar(slide, spec: dict[str, Any], theme: Theme):
    main_block = spec["blocks"][0] if spec["blocks"] else {"title": spec["title"], "items": [], "summary": ""}
    add_card(slide, x=0.72, y=1.62, w=8.22, h=4.9, theme=theme, accent=theme.primary, title=main_block["title"], items=main_block["items"], summary=main_block["summary"] or spec.get("page_goal", ""))

    sidebar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(9.18),
        Inches(1.62),
        Inches(3.38),
        Inches(4.9),
    )
    sidebar.fill.solid()
    sidebar.fill.fore_color.rgb = rgb(theme.surface, "F5F7FA")
    sidebar.line.color.rgb = rgb(theme.divider, "E5E7EB")
    sidebar.line.width = Pt(1.0)
    add_textbox(
        slide,
        x=9.36,
        y=1.78,
        w=3.0,
        h=0.22,
        text=spec.get("visual_focus") or "视觉焦点",
        font_face=theme.body_font,
        font_size=13,
        color=theme.accent,
        bold=True,
    )
    sidebar_items = split_summary(spec.get("detail_notes") or spec.get("page_goal") or "突出关键结构与价值主张。", 16)
    add_bullets(slide, x=9.36, y=2.12, w=2.92, h=3.7, items=sidebar_items, theme=theme, font_size=11)


def render_slide(slide, spec: dict[str, Any], theme: Theme):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(theme.background, "FFFFFF")
    add_header(slide, spec, theme)
    family = spec.get("layout_family", "focus_with_sidebar")
    if family == "architecture_layers":
        render_architecture_layers(slide, spec, theme)
    elif family == "roadmap":
        render_roadmap(slide, spec, theme)
    elif family == "two_column":
        render_two_column(slide, spec, theme)
    elif family == "four_card_grid":
        render_card_grid(slide, spec, theme, columns=2)
    elif family == "three_card_grid":
        render_card_grid(slide, spec, theme, columns=3)
    elif family == "summary":
        render_summary(slide, spec, theme)
    else:
        render_focus_with_sidebar(slide, spec, theme)


def render_ppt_from_specs(slide_specs: dict[str, Any], master_style: dict[str, Any], output_path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    theme = to_theme(master_style)
    blank = prs.slide_layouts[6]
    for spec in slide_specs.get("slides", []):
        slide = prs.slides.add_slide(blank)
        render_slide(slide, spec, theme)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path
