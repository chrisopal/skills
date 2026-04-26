"""Render a PPTX directly from slide_prompts.json (pattern + slots) and master_style.json.

Bypasses the legacy LLM-driven JS path entirely: dispatches each
pattern-mode slide to a deterministic Python renderer in
`scripts/lib/pattern_pptx_renderers.py`. Slides without a pattern_id
or with layout_mode != "pattern" fall back to a generic title-only
slide so the deck still produces output.

Usage:
    python scripts/render_pptx_from_intents.py
        --slide-prompts artifacts/slide_prompts.json
        --master-style artifacts/master_style.json
        --output deck.pptx
"""

from __future__ import annotations

import argparse
import importlib.util
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

SKILL_ROOT = Path(__file__).resolve().parent.parent
PATTERNS_DIR = SKILL_ROOT / "assets" / "patterns"

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5


def _load(name: str, rel_path: str):
    if name in sys.modules:
        return sys.modules[name]
    spec = importlib.util.spec_from_file_location(name, SKILL_ROOT / rel_path)
    module = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    sys.modules[name] = module
    spec.loader.exec_module(module)
    return module


renderers_mod = _load("pattern_pptx_renderers", "scripts/lib/pattern_pptx_renderers.py")
registry_mod = _load("pattern_registry", "scripts/lib/pattern_registry.py")


def _flatten_slots(slots: dict) -> dict:
    """Page intents may nest slots like cell_1: {label, value, desc}."""

    out: dict[str, str] = {}
    for key, value in (slots or {}).items():
        if isinstance(value, dict):
            for sub_key, sub_value in value.items():
                out[f"{key}_{sub_key}"] = "" if sub_value is None else str(sub_value)
        else:
            out[key] = "" if value is None else str(value)
    return out


def render_deck(
    slide_prompts: dict,
    master_style: dict,
    output_path: Path,
    *,
    patterns_dir: Path = PATTERNS_DIR,
) -> Path:
    pres = Presentation()
    pres.slide_width = Inches(SLIDE_W_IN)
    pres.slide_height = Inches(SLIDE_H_IN)
    blank_layout = pres.slide_layouts[6]

    registry = registry_mod.PatternRegistry(patterns_dir)
    theme = renderers_mod.theme_from_master_style(master_style)

    for slide_data in slide_prompts.get("slides", []):
        slide = pres.slides.add_slide(blank_layout)
        pattern_id = slide_data.get("pattern_id")
        layout_mode = slide_data.get("layout_mode", "custom")
        slots = _flatten_slots(slide_data.get("slots") or {})
        if pattern_id and layout_mode == "pattern" and pattern_id in registry:
            regions = registry.get(pattern_id).layout_regions
            renderers_mod.render_slide(slide, pattern_id, slots, regions, theme)
        else:
            # Fallback: just put the page title
            renderers_mod._add_text(
                slide, slide_data.get("title", ""),
                0.5, 0.5, 12.3, 1.0,
                size=28, bold=True, color=theme["text_primary"],
            )
            renderers_mod._add_text(
                slide, "[ no pattern_id — supply pattern + slots for a styled render ]",
                0.5, 6.6, 12.3, 0.4,
                size=11, color=theme["divider"],
            )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    pres.save(str(output_path))
    return output_path


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Render PPTX directly from slide_prompts patterns.")
    parser.add_argument("--slide-prompts", required=True, type=Path)
    parser.add_argument("--master-style", required=True, type=Path)
    parser.add_argument("--output", required=True, type=Path)
    args = parser.parse_args(argv)

    slide_prompts = json.loads(args.slide_prompts.read_text(encoding="utf-8"))
    master_style = json.loads(args.master_style.read_text(encoding="utf-8"))
    out = render_deck(slide_prompts, master_style, args.output)
    print(f"[OK] Wrote PPTX: {out}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
