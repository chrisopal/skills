"""Verify each Python pattern renderer produces shapes containing the slot data."""

from __future__ import annotations

import importlib.util
import json
import sys
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

SKILL_ROOT = Path(__file__).resolve().parent.parent


def _load(name: str, rel_path: str):
    if name in sys.modules:
        return sys.modules[name]
    spec = importlib.util.spec_from_file_location(name, SKILL_ROOT / rel_path)
    module = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    sys.modules[name] = module
    spec.loader.exec_module(module)
    return module


renderers = _load("pattern_pptx_renderers", "scripts/lib/pattern_pptx_renderers.py")
driver = _load("render_pptx_from_intents", "scripts/render_pptx_from_intents.py")


@pytest.fixture
def fresh_slide():
    pres = Presentation()
    pres.slide_width = Inches(13.333)
    pres.slide_height = Inches(7.5)
    return pres.slides.add_slide(pres.slide_layouts[6])


@pytest.fixture
def theme():
    return renderers.theme_from_master_style(None)


def _slide_text(slide) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    parts.append(run.text)
    return "\n".join(parts)


def test_theme_from_master_style_handles_huixin_keys():
    style = {"color_strategy": {
        "primary_green": "#A8D86B",
        "secondary_teal": "#0F95B6",
        "text_primary": "#1E1E1E",
    }}
    theme = renderers.theme_from_master_style(style)
    assert theme["primary"] == "#A8D86B"
    assert theme["secondary"] == "#0F95B6"
    assert theme["text_primary"] == "#1E1E1E"


def test_theme_from_master_style_handles_dark_keys():
    style = {"color_strategy": {
        "primary_blue": "#38BDF8",
        "background": "#0B1020",
    }}
    theme = renderers.theme_from_master_style(style)
    assert theme["primary"] == "#38BDF8"
    assert theme["background"] == "#0B1020"


def test_render_cover_includes_title_and_org(fresh_slide, theme):
    slots = {
        "title": "FY26 Q1",
        "subtitle": "高管复盘",
        "org_block": "ACME 集团",
    }
    renderers.render_cover(fresh_slide, slots, {}, theme)
    text = _slide_text(fresh_slide)
    assert "FY26 Q1" in text
    assert "ACME 集团" in text


def test_render_kpi_strip_shows_all_provided_values(fresh_slide, theme):
    slots = {
        "kpi_1_value": "+38%", "kpi_1_label": "ROI",
        "kpi_2_value": "12k", "kpi_2_label": "MAU",
        "kpi_3_value": "62", "kpi_3_label": "NPS",
        "kpi_4_value": "100M", "kpi_4_label": "GMV",
    }
    renderers.render_kpi_strip(fresh_slide, slots, {}, theme)
    text = _slide_text(fresh_slide)
    for needle in ("+38%", "12k", "62", "100M", "ROI", "MAU", "NPS", "GMV"):
        assert needle in text, f"missing {needle} in rendered slide"


def test_render_kpi_strip_skips_empty_cells(fresh_slide, theme):
    slots = {
        "kpi_1_value": "+38%", "kpi_1_label": "ROI",
        "kpi_2_value": "12k", "kpi_2_label": "MAU",
    }
    renderers.render_kpi_strip(fresh_slide, slots, {}, theme)
    text = _slide_text(fresh_slide)
    assert "+38%" in text
    assert "12k" in text
    # No phantom cells
    assert "kpi_3" not in text


def test_render_summary_takeaways_numbers_each_takeaway(fresh_slide, theme):
    slots = {
        "takeaway_1": "Action A",
        "takeaway_2": "Action B",
        "takeaway_3": "Action C",
        "next_step": "Q2 review",
    }
    renderers.render_summary_takeaways(fresh_slide, slots, {}, theme)
    text = _slide_text(fresh_slide)
    assert "Action A" in text
    assert "Action B" in text
    assert "Action C" in text
    assert "Q2 review" in text
    # Numbered badges
    for digit in ("1", "2", "3"):
        assert digit in text


def test_render_four_card_matrix_includes_value_and_desc(fresh_slide, theme):
    slots = {
        "cell_1_label": "Eff", "cell_1_value": "+38%", "cell_1_desc": "demo",
        "cell_2_label": "Cost", "cell_2_value": "-12%", "cell_2_desc": "demo",
        "cell_3_label": "MAU", "cell_3_value": "12k", "cell_3_desc": "demo",
        "cell_4_label": "NPS", "cell_4_value": "62", "cell_4_desc": "demo",
    }
    renderers.render_four_card_matrix(fresh_slide, slots, {}, theme)
    text = _slide_text(fresh_slide)
    assert "+38%" in text
    assert "-12%" in text
    assert "Eff" in text


def test_render_two_column_compare_includes_verdict(fresh_slide, theme):
    slots = {
        "left_title": "Old", "left_point_1": "p1",
        "right_title": "New", "right_point_1": "q1",
        "verdict": "New wins",
    }
    renderers.render_two_column_compare(fresh_slide, slots, {}, theme)
    text = _slide_text(fresh_slide)
    assert "Old" in text and "New" in text
    assert "New wins" in text


def test_render_section_divider_renders_section_no_and_title(fresh_slide, theme):
    slots = {"section_no": "02", "section_title": "执行计划", "agenda_anchor": "Q2 行动"}
    renderers.render_section_divider(fresh_slide, slots, {}, theme)
    text = _slide_text(fresh_slide)
    assert "02" in text and "执行计划" in text and "Q2 行动" in text


def test_render_freeform_falls_back_to_headline_and_note(fresh_slide, theme):
    slots = {"headline": "Custom", "note": "Need design"}
    renderers.render_freeform(fresh_slide, slots, {}, theme)
    text = _slide_text(fresh_slide)
    assert "Custom" in text and "Need design" in text


def test_render_slide_dispatches_known_pattern(fresh_slide, theme):
    assert renderers.render_slide(
        fresh_slide, "cover", {"title": "X"}, {}, theme,
    ) is True


def test_render_slide_returns_false_for_unknown_pattern(fresh_slide, theme):
    assert renderers.render_slide(
        fresh_slide, "made_up_pattern", {}, {}, theme,
    ) is False


def test_all_twelve_patterns_have_renderer():
    expected = {
        "cover", "section_divider", "conclusion_top_modules", "two_column_compare",
        "four_card_matrix", "three_stage_path", "kpi_strip", "architecture_layers",
        "before_after", "evidence_grid", "summary_takeaways", "freeform",
    }
    assert set(renderers.PATTERN_RENDERERS.keys()) == expected


# ---------------------------------------------------------------------------
# integration via render_pptx_from_intents
# ---------------------------------------------------------------------------


def test_render_deck_produces_pptx_with_slot_data(tmp_path):
    slide_prompts = {
        "slides": [
            {
                "page_no": 1, "title": "封面", "page_goal": "demo",
                "layout_type": "cover", "key_blocks": [], "compiled_prompt": "x",
                "intent_status": "locked",
                "pattern_id": "cover", "layout_mode": "pattern",
                "slots": {"title": "Demo Cover", "subtitle": "Sub", "org_block": "Org"},
            },
            {
                "page_no": 2, "title": "KPI", "page_goal": "demo",
                "layout_type": "kpi_strip", "key_blocks": [], "compiled_prompt": "x",
                "intent_status": "locked",
                "pattern_id": "kpi_strip", "layout_mode": "pattern",
                "slots": {
                    "kpi_1_value": "+38%", "kpi_1_label": "ROI",
                    "kpi_2_value": "12k", "kpi_2_label": "MAU",
                },
            },
        ],
        "quality_checklist": {},
    }
    master_style = {
        "color_strategy": {
            "primary_green": "#A8D86B",
            "secondary_teal": "#0F95B6",
            "background": "#FFFFFF",
            "text_primary": "#1E1E1E",
        }
    }
    out = tmp_path / "deck.pptx"
    driver.render_deck(slide_prompts, master_style, out)
    assert out.exists()

    pres = Presentation(str(out))
    assert len(pres.slides) == 2
    text_slide_1 = "\n".join(
        run.text for shape in pres.slides[0].shapes if shape.has_text_frame
        for para in shape.text_frame.paragraphs for run in para.runs
    )
    assert "Demo Cover" in text_slide_1
    text_slide_2 = "\n".join(
        run.text for shape in pres.slides[1].shapes if shape.has_text_frame
        for para in shape.text_frame.paragraphs for run in para.runs
    )
    assert "+38%" in text_slide_2
    assert "ROI" in text_slide_2


def test_render_deck_falls_back_for_slides_without_pattern(tmp_path):
    slide_prompts = {
        "slides": [
            {
                "page_no": 1, "title": "Free slide", "page_goal": "demo",
                "layout_type": "custom", "key_blocks": [], "compiled_prompt": "x",
                "intent_status": "locked", "layout_mode": "custom",
            }
        ],
        "quality_checklist": {},
    }
    master_style = {"color_strategy": {"primary_green": "#A8D86B"}}
    out = tmp_path / "deck.pptx"
    driver.render_deck(slide_prompts, master_style, out)
    pres = Presentation(str(out))
    assert len(pres.slides) == 1
    text = "\n".join(
        run.text for shape in pres.slides[0].shapes if shape.has_text_frame
        for para in shape.text_frame.paragraphs for run in para.runs
    )
    assert "Free slide" in text
