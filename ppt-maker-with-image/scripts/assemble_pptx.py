#!/usr/bin/env python3
from __future__ import annotations

import argparse
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches


SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Assemble slide images into a PPTX file.")
    parser.add_argument("--images", nargs="+", required=True, help="Ordered list of slide image paths")
    parser.add_argument("--output", required=True, help="Output .pptx path")
    return parser


def validate_images(image_paths: list[Path]) -> None:
    for path in image_paths:
        if not path.exists():
            raise FileNotFoundError(f"Image not found: {path}")
        with Image.open(path) as image:
            image.verify()

def assemble_pptx(image_paths: list[Path], output_path: Path) -> Path:
    validate_images(image_paths)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    blank_layout = prs.slide_layouts[6]

    for image_path in image_paths:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(image_path),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    if prs.slides:
        xml_slides = prs.slides._sldIdLst  # type: ignore[attr-defined]
        first = xml_slides[0]
        xml_slides.remove(first)

    prs.save(str(output_path))
    return output_path


def main() -> int:
    parser = build_parser()
    args = parser.parse_args()

    image_paths = [Path(p).expanduser().resolve() for p in args.images]
    output_path = Path(args.output).expanduser().resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)

    assemble_pptx(image_paths, output_path)
    print(f"[OK] Wrote PPTX: {output_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
