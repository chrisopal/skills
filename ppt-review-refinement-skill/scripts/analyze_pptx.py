#!/usr/bin/env python3
"""Structural and quantitative audit for an existing PPTX deck."""

from __future__ import annotations

import argparse
import io
import sys
from collections import Counter
from datetime import datetime, timezone
from pathlib import Path
from typing import Any
from zipfile import BadZipFile, ZipFile

from PIL import Image, UnidentifiedImageError
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from common import (
    bbox_area,
    color_descriptor,
    emu_to_in,
    emu_to_pt,
    ensure_dir,
    explicit_color_hex,
    intersection_area,
    is_background_like,
    iter_nested_shapes,
    iter_runs,
    load_yaml,
    median,
    safe_fill_hex,
    safe_line_hex,
    safe_shape_name,
    safe_shape_type_name,
    sha256_bytes,
    sha256_file,
    shape_bbox,
    shape_text,
    slide_title_shape,
    text_overflow_estimate,
    text_by_slide,
    write_json,
)

SKILL_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_CONFIG = SKILL_ROOT / "config" / "defaults.yaml"


def make_issue(
    issue_id: str,
    severity: str,
    category: str,
    scope: str,
    finding: str,
    recommendation: str,
    *,
    slide: int | str = "global",
    evidence: list[str] | None = None,
    confidence: str = "high",
) -> dict[str, Any]:
    return {
        "id": issue_id,
        "severity": severity,
        "category": category,
        "scope": scope,
        "slide": slide,
        "finding": finding,
        "recommendation": recommendation,
        "evidence": evidence or [],
        "confidence": confidence,
    }


def package_features(path: Path) -> dict[str, Any]:
    features = {
        "has_vba": False,
        "diagram_parts": 0,
        "embedded_objects": 0,
        "slide_timings": 0,
        "external_relationships": 0,
        "hyperlink_markers": 0,
        "notes_slides": 0,
        "media_files": 0,
    }
    try:
        with ZipFile(path, "r") as zf:
            names = zf.namelist()
            features["has_vba"] = "ppt/vbaProject.bin" in names
            features["diagram_parts"] = sum(name.startswith("ppt/diagrams/") for name in names)
            features["embedded_objects"] = sum(name.startswith("ppt/embeddings/") for name in names)
            features["notes_slides"] = sum(
                name.startswith("ppt/notesSlides/notesSlide") and name.endswith(".xml") for name in names
            )
            features["media_files"] = sum(name.startswith("ppt/media/") for name in names)
            for name in names:
                if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                    data = zf.read(name)
                    features["slide_timings"] += data.count(b"<p:timing")
                    features["hyperlink_markers"] += data.count(b"hlinkClick")
                if name.endswith(".rels"):
                    data = zf.read(name)
                    features["external_relationships"] += data.count(b'TargetMode="External"')
    except BadZipFile:
        pass
    return features


def image_info(shape: Any) -> dict[str, Any] | None:
    try:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            return None
        image = shape.image
        blob = image.blob
        ext = image.ext
        content_type = image.content_type
    except (AttributeError, ValueError, KeyError):
        return None

    display_w = emu_to_in(shape.width)
    display_h = emu_to_in(shape.height)
    info: dict[str, Any] = {
        "shape_id": getattr(shape, "shape_id", None),
        "name": safe_shape_name(shape),
        "extension": ext,
        "content_type": content_type,
        "display_width_in": display_w,
        "display_height_in": display_h,
        "sha256": sha256_bytes(blob),
        "pixel_width": None,
        "pixel_height": None,
        "effective_dpi": None,
        "raster_readable": False,
    }
    try:
        with Image.open(io.BytesIO(blob)) as im:
            px_w, px_h = im.size
        info["pixel_width"] = px_w
        info["pixel_height"] = px_h
        info["raster_readable"] = True
        if display_w > 0 and display_h > 0:
            info["effective_dpi"] = round(min(px_w / display_w, px_h / display_h), 1)
    except (UnidentifiedImageError, OSError, ValueError):
        # Vector formats and some embedded encodings cannot be read by Pillow.
        pass
    return info


def analyze_deck(input_path: str | Path, config_path: str | Path | None = None) -> dict[str, Any]:
    input_path = Path(input_path).resolve()
    if not input_path.exists():
        raise FileNotFoundError(input_path)
    if input_path.suffix.lower() not in {".pptx", ".pptm"}:
        raise ValueError("Only .pptx/.pptm files are structurally supported. Convert .ppt first.")

    config = load_yaml(config_path or DEFAULT_CONFIG)
    thresholds = config.get("review", {})
    prs = Presentation(str(input_path))
    slide_width_in = emu_to_in(prs.slide_width)
    slide_height_in = emu_to_in(prs.slide_height)

    global_fonts: Counter[str] = Counter()
    global_sizes: Counter[str] = Counter()
    global_font_colors: Counter[str] = Counter()
    global_fill_colors: Counter[str] = Counter()
    global_line_colors: Counter[str] = Counter()
    global_shape_types: Counter[str] = Counter()
    global_layouts: Counter[str] = Counter()
    all_issues: list[dict[str, Any]] = []
    slides_out: list[dict[str, Any]] = []
    title_x_values: list[tuple[int, float]] = []
    title_y_values: list[tuple[int, float]] = []
    title_size_values: list[tuple[int, float]] = []
    issue_counter = 1

    def next_id(prefix: str) -> str:
        nonlocal issue_counter
        value = f"{prefix}-{issue_counter:04d}"
        issue_counter += 1
        return value

    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_issues: list[dict[str, Any]] = []
        shape_type_counter: Counter[str] = Counter()
        fonts: Counter[str] = Counter()
        sizes: Counter[str] = Counter()
        font_colors: Counter[str] = Counter()
        fill_colors: Counter[str] = Counter()
        line_colors: Counter[str] = Counter()
        shape_records: list[dict[str, Any]] = []
        images: list[dict[str, Any]] = []
        small_text_runs: list[dict[str, Any]] = []
        top_level_geometry: list[tuple[Any, tuple[float, float, float, float]]] = []

        try:
            layout_name = slide.slide_layout.name or "unnamed"
        except (AttributeError, ValueError):
            layout_name = "unknown"
        global_layouts[layout_name] += 1

        title_shape = slide_title_shape(slide)
        title_shape_id = getattr(title_shape, "shape_id", None)
        title_text = shape_text(title_shape).strip() if title_shape is not None else ""
        title_box = shape_bbox(title_shape) if title_shape is not None else None
        if title_box and slide_index > 1:
            title_x_values.append((slide_index, title_box[0]))
            title_y_values.append((slide_index, title_box[1]))

        for shape in iter_nested_shapes(slide.shapes):
            type_name = safe_shape_type_name(shape)
            shape_type_counter[type_name] += 1
            global_shape_types[type_name] += 1

            fill_hex = safe_fill_hex(shape)
            if fill_hex:
                fill_colors[fill_hex] += 1
                global_fill_colors[fill_hex] += 1
            line_hex = safe_line_hex(shape)
            if line_hex:
                line_colors[line_hex] += 1
                global_line_colors[line_hex] += 1

            for _tf, paragraph, run in iter_runs(shape):
                text = run.text or ""
                if not text.strip():
                    continue
                font_name = run.font.name or "<theme/default>"
                size_pt = emu_to_pt(run.font.size)
                size_label = f"{size_pt:g}pt" if size_pt is not None else "<theme/default>"
                color_label = color_descriptor(run.font.color)
                weight = max(1, len(text.strip()))
                fonts[font_name] += weight
                sizes[size_label] += weight
                font_colors[color_label] += weight
                global_fonts[font_name] += weight
                global_sizes[size_label] += weight
                global_font_colors[color_label] += weight

                if size_pt is not None and size_pt < float(thresholds.get("min_note_font_pt", 9)):
                    small_text_runs.append(
                        {
                            "shape_id": getattr(shape, "shape_id", None),
                            "shape_name": safe_shape_name(shape),
                            "text": text[:80],
                            "size_pt": size_pt,
                            "paragraph_level": getattr(paragraph, "level", 0),
                        }
                    )
                if getattr(shape, "shape_id", None) == title_shape_id and size_pt is not None:
                    title_size_values.append((slide_index, size_pt))

            overflow = text_overflow_estimate(shape)
            if overflow and overflow["ratio"] >= float(thresholds.get("text_overflow_ratio_warning", 1.25)):
                critical_ratio = float(thresholds.get("text_overflow_ratio_critical", 2.0))
                severity = "P0" if overflow["ratio"] >= critical_ratio else "P1"
                issue = make_issue(
                    next_id("TXT-OVERFLOW"),
                    severity,
                    "typography",
                    "object",
                    f"文本框的估算文字量约为可用高度的 {overflow['ratio']:.1f} 倍，存在明显溢出风险。",
                    "扩大文本框或减少内容；不要仅通过继续缩小字号来掩盖溢出。",
                    slide=slide_index,
                    evidence=[
                        f"estimated_lines={overflow['estimated_lines']:.0f}",
                        f"available_lines={overflow['available_lines']:.2f}",
                        f"font_size={overflow['font_size_pt']:.1f}pt",
                    ],
                    confidence="medium",
                )
                slide_issues.append(issue)
                all_issues.append(issue)

            info = image_info(shape)
            if info:
                images.append(info)
                dpi = info.get("effective_dpi")
                if dpi is not None and dpi < float(thresholds.get("image_dpi_warning", 120)):
                    severity = "P0" if dpi < float(thresholds.get("image_dpi_critical", 72)) else "P1"
                    issue = make_issue(
                        next_id("IMG"),
                        severity,
                        "imagery",
                        "object",
                        f"图片有效分辨率约为 {dpi} DPI，放映或导出时可能模糊。",
                        "替换为更高分辨率源图，或降低该图片在页面中的显示尺寸。",
                        slide=slide_index,
                        evidence=[info.get("name", "picture")],
                    )
                    slide_issues.append(issue)
                    all_issues.append(issue)

        # Top-level geometry is used for overlap/out-of-bounds checks to avoid group self-overlap noise.
        for shape in slide.shapes:
            box = shape_bbox(shape)
            if box is None:
                continue
            type_name = safe_shape_type_name(shape)
            record = {
                "shape_id": getattr(shape, "shape_id", None),
                "name": safe_shape_name(shape),
                "type": type_name,
                "bbox_in": {"x": box[0], "y": box[1], "w": box[2], "h": box[3]},
                "text_preview": shape_text(shape).replace("\n", " ")[:140],
            }
            shape_records.append(record)

            tolerance = float(thresholds.get("edge_tolerance_in", 0.03))
            x, y, w, h = box
            if (
                x < -tolerance
                or y < -tolerance
                or x + w > slide_width_in + tolerance
                or y + h > slide_height_in + tolerance
            ):
                issue = make_issue(
                    next_id("GEO"),
                    "P0",
                    "layout",
                    "object",
                    f"对象“{record['name']}”超出页面边界。",
                    "将对象移回安全区域，或确认这是有意的出血背景。",
                    slide=slide_index,
                    evidence=[str(record["bbox_in"])],
                    confidence="medium" if is_background_like(box, slide_width_in, slide_height_in) else "high",
                )
                slide_issues.append(issue)
                all_issues.append(issue)

            if type_name not in {"LINE", "CONNECTOR"} and not is_background_like(
                box, slide_width_in, slide_height_in
            ):
                top_level_geometry.append((shape, box))

        overlap_threshold = float(thresholds.get("overlap_area_ratio_warning", 0.08))
        for i in range(len(top_level_geometry)):
            shape_a, box_a = top_level_geometry[i]
            area_a = bbox_area(box_a)
            if area_a < 0.01:
                continue
            for j in range(i + 1, len(top_level_geometry)):
                shape_b, box_b = top_level_geometry[j]
                area_b = bbox_area(box_b)
                if area_b < 0.01:
                    continue
                inter = intersection_area(box_a, box_b)
                if inter <= 0:
                    continue
                ratio = inter / min(area_a, area_b)
                # Full containment is usually intentional (text/image placed on a card/background).
                if ratio >= 0.92:
                    continue
                if ratio >= overlap_threshold:
                    issue = make_issue(
                        next_id("OVR"),
                        "P1",
                        "layout",
                        "object",
                        f"对象“{safe_shape_name(shape_a)}”与“{safe_shape_name(shape_b)}”存在明显交叠风险。",
                        "检查渲染图确认是否为有意叠放；若影响阅读，调整位置、层级或间距。",
                        slide=slide_index,
                        evidence=[f"overlap_ratio={ratio:.2f}"],
                        confidence="medium",
                    )
                    slide_issues.append(issue)
                    all_issues.append(issue)

        if small_text_runs:
            issue = make_issue(
                next_id("TXT"),
                "P1",
                "typography",
                "slide",
                f"检测到 {len(small_text_runs)} 处小于 {thresholds.get('min_note_font_pt', 9)}pt 的文字。",
                "确认是否仅为脚注；正文或图表标签应提高字号或减少信息密度。",
                slide=slide_index,
                evidence=[f"{item['size_pt']}pt: {item['text']}" for item in small_text_runs[:5]],
                confidence="medium",
            )
            slide_issues.append(issue)
            all_issues.append(issue)

        texts = text_by_slide(slide)
        char_count = len(texts["all"].replace("\n", ""))
        dense_threshold = int(thresholds.get("dense_slide_character_count", 420))
        very_dense_threshold = int(thresholds.get("very_dense_slide_character_count", 700))
        if char_count > dense_threshold:
            severity = "P1" if char_count <= very_dense_threshold else "P0"
            issue = make_issue(
                next_id("DNS"),
                severity,
                "hierarchy",
                "slide",
                f"页面文本约 {char_count} 个字符，信息密度偏高。",
                "提炼结论、减少同权重信息；必要时拆分页面，但拆分需获得授权。",
                slide=slide_index,
                evidence=[f"characters={char_count}"],
            )
            slide_issues.append(issue)
            all_issues.append(issue)

        slides_out.append(
            {
                "slide": slide_index,
                "layout_name": layout_name,
                "title": title_text,
                "title_box_in": (
                    {"x": title_box[0], "y": title_box[1], "w": title_box[2], "h": title_box[3]}
                    if title_box
                    else None
                ),
                "text": texts,
                "character_count": char_count,
                "shape_count_top_level": len(slide.shapes),
                "shape_types": dict(shape_type_counter.most_common()),
                "fonts_weighted_by_characters": dict(fonts.most_common()),
                "font_sizes_weighted_by_characters": dict(sizes.most_common()),
                "font_colors_weighted_by_characters": dict(font_colors.most_common()),
                "fill_colors": dict(fill_colors.most_common()),
                "line_colors": dict(line_colors.most_common()),
                "images": images,
                "shapes": shape_records,
                "issues": slide_issues,
            }
        )

    explicit_fonts = [name for name in global_fonts if name != "<theme/default>"]
    if len(explicit_fonts) > int(thresholds.get("max_recommended_fonts", 3)):
        all_issues.append(
            make_issue(
                next_id("GLB"),
                "P1",
                "typography",
                "deck",
                f"全稿检测到 {len(explicit_fonts)} 种显式字体，字体体系较分散。",
                "建立中文/英文字体配对，并按封面、标题、正文、数据、注释等角色统一。",
                evidence=explicit_fonts[:12],
            )
        )

    explicit_sizes = [size for size in global_sizes if size != "<theme/default>"]
    if len(explicit_sizes) > int(thresholds.get("max_recommended_font_sizes", 10)):
        all_issues.append(
            make_issue(
                next_id("GLB"),
                "P1",
                "typography",
                "deck",
                f"全稿检测到 {len(explicit_sizes)} 种显式字号，层级标尺不稳定。",
                "将字号收敛为有限角色等级，避免每页凭感觉设置。",
                evidence=explicit_sizes[:15],
            )
        )

    explicit_colors = set(global_fill_colors) | set(global_line_colors) | {
        c for c in global_font_colors if c.startswith("#")
    }
    if len(explicit_colors) > int(thresholds.get("max_recommended_explicit_colors", 14)):
        all_issues.append(
            make_issue(
                next_id("GLB"),
                "P1",
                "color",
                "deck",
                f"全稿检测到 {len(explicit_colors)} 种显式 RGB 颜色，颜色角色可能不清。",
                "先区分品牌色、强调色和语义色，再执行精确颜色映射。",
                evidence=sorted(explicit_colors)[:20],
            )
        )

    title_x_median = median([v for _, v in title_x_values])
    title_y_median = median([v for _, v in title_y_values])
    title_tolerance = float(thresholds.get("title_position_tolerance_in", 0.18))
    title_position_outliers: list[str] = []
    if title_x_median is not None and title_y_median is not None:
        x_map = dict(title_x_values)
        y_map = dict(title_y_values)
        for slide_num in sorted(set(x_map) & set(y_map)):
            if abs(x_map[slide_num] - title_x_median) > title_tolerance or abs(
                y_map[slide_num] - title_y_median
            ) > title_tolerance:
                title_position_outliers.append(
                    f"slide {slide_num}: x={x_map[slide_num]:.2f}, y={y_map[slide_num]:.2f}"
                )
    if title_position_outliers:
        all_issues.append(
            make_issue(
                next_id("GLB"),
                "P1",
                "layout",
                "deck",
                f"有 {len(title_position_outliers)} 页标题位置明显偏离主流位置。",
                "按标题区 Token 统一左边界、顶部、宽度和高度；封面与章节页可单独处理。",
                evidence=title_position_outliers[:12],
                confidence="medium",
            )
        )

    unique_title_sizes = sorted({round(v, 1) for _, v in title_size_values})
    if len(unique_title_sizes) > 4:
        all_issues.append(
            make_issue(
                next_id("GLB"),
                "P1",
                "typography",
                "deck",
                f"标题文本检测到 {len(unique_title_sizes)} 种字号。",
                "区分封面、章节和内容页标题，统一每类标题字号。",
                evidence=[f"{v:g}pt" for v in unique_title_sizes],
                confidence="medium",
            )
        )

    features = package_features(input_path)
    limitations: list[str] = []
    if features["diagram_parts"]:
        limitations.append("检测到 SmartArt/Diagram 部件；结构脚本不能完整判断其可编辑和渲染状态。")
    if features["embedded_objects"]:
        limitations.append("检测到嵌入对象；需要在 PowerPoint 中人工检查。")
    if features["has_vba"]:
        limitations.append("检测到 VBA；不要用不支持宏的编辑器覆盖保存。")
    if features["slide_timings"]:
        limitations.append("检测到动画/时间线；python-pptx 不会完整保留或编辑动画。")
    limitations.append("对象重叠是几何风险提示，必须结合逐页渲染图判断是否为有意叠放。")
    limitations.append("继承自主题/母版的字体、颜色和字号可能显示为 <theme/default>，需结合渲染结果。")
    limitations.append("文本溢出检测是基于字号、文本量和文本框尺寸的启发式估算，仍需结合逐页渲染图确认。")

    aspect = round(slide_width_in / slide_height_in, 4) if slide_height_in else None
    return {
        "version": "1.0.0",
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "source": {
            "path": str(input_path),
            "filename": input_path.name,
            "sha256": sha256_file(input_path),
            "size_bytes": input_path.stat().st_size,
        },
        "deck": {
            "slide_count": len(prs.slides),
            "slide_width_in": slide_width_in,
            "slide_height_in": slide_height_in,
            "aspect_ratio": aspect,
            "slide_master_count": len(prs.slide_masters),
            "slide_layout_count": len(prs.slide_layouts),
            "package_features": features,
        },
        "global_stats": {
            "layouts": dict(global_layouts.most_common()),
            "shape_types": dict(global_shape_types.most_common()),
            "fonts_weighted_by_characters": dict(global_fonts.most_common()),
            "font_sizes_weighted_by_characters": dict(global_sizes.most_common()),
            "font_colors_weighted_by_characters": dict(global_font_colors.most_common()),
            "fill_colors": dict(global_fill_colors.most_common()),
            "line_colors": dict(global_line_colors.most_common()),
            "explicit_rgb_color_count": len(explicit_colors),
            "title_position_median_in": {"x": title_x_median, "y": title_y_median},
            "title_size_variants_pt": unique_title_sizes,
        },
        "slides": slides_out,
        "issues": all_issues,
        "limitations": limitations,
    }


def markdown_report(analysis: dict[str, Any]) -> str:
    deck = analysis["deck"]
    stats = analysis["global_stats"]
    issues = analysis["issues"]
    lines = [
        "# PPTX 结构审计",
        "",
        f"- 文件：`{analysis['source']['filename']}`",
        f"- SHA-256：`{analysis['source']['sha256']}`",
        f"- 页数：{deck['slide_count']}",
        f"- 尺寸：{deck['slide_width_in']} × {deck['slide_height_in']} in",
        f"- 宽高比：{deck['aspect_ratio']}",
        f"- 母版：{deck['slide_master_count']}；版式：{deck['slide_layout_count']}",
        "",
        "## 全局统计",
        "",
        f"- 字体：{', '.join(list(stats['fonts_weighted_by_characters'])[:10]) or '未检测到显式字体'}",
        f"- 字号：{', '.join(list(stats['font_sizes_weighted_by_characters'])[:12]) or '未检测到显式字号'}",
        f"- 显式 RGB 颜色数量：{stats['explicit_rgb_color_count']}",
        f"- 主要版式：{', '.join(f'{k}({v})' for k, v in list(stats['layouts'].items())[:8])}",
        "",
        "## 问题概览",
        "",
    ]
    severity_counts = Counter(issue["severity"] for issue in issues)
    lines.append(
        f"- P0：{severity_counts.get('P0', 0)}；P1：{severity_counts.get('P1', 0)}；P2：{severity_counts.get('P2', 0)}"
    )
    lines.extend(["", "| 级别 | 页码 | 类别 | 发现 | 建议 |", "|---|---:|---|---|---|"])
    order = {"P0": 0, "P1": 1, "P2": 2}
    for issue in sorted(issues, key=lambda x: (order.get(x["severity"], 9), str(x["slide"]), x["id"])):
        finding = issue["finding"].replace("|", "\\|").replace("\n", " ")
        recommendation = issue["recommendation"].replace("|", "\\|").replace("\n", " ")
        lines.append(
            f"| {issue['severity']} | {issue['slide']} | {issue['category']} | {finding} | {recommendation} |"
        )

    lines.extend(["", "## 逐页摘要", "", "| 页码 | 标题 | 字符数 | 对象数 | 图片数 | 问题数 |", "|---:|---|---:|---:|---:|---:|"])
    for slide in analysis["slides"]:
        title = (slide["title"] or "（无标题）").replace("|", "\\|").replace("\n", " ")[:60]
        lines.append(
            f"| {slide['slide']} | {title} | {slide['character_count']} | {slide['shape_count_top_level']} | {len(slide['images'])} | {len(slide['issues'])} |"
        )

    lines.extend(["", "## 需要人工复核", ""])
    for item in analysis.get("limitations", []):
        lines.append(f"- {item}")
    lines.append("")
    return "\n".join(lines)


def main() -> int:
    parser = argparse.ArgumentParser(description="Audit PPTX structure, typography, colors, images, and geometry.")
    parser.add_argument("input", help="Input .pptx/.pptm file")
    parser.add_argument("--out", required=True, help="Output directory")
    parser.add_argument("--config", help="Optional defaults YAML")
    args = parser.parse_args()

    try:
        out_dir = ensure_dir(args.out)
        analysis = analyze_deck(args.input, args.config)
        write_json(out_dir / "analysis.json", analysis)
        (out_dir / "structural_audit.md").write_text(markdown_report(analysis), encoding="utf-8")
        print(f"Wrote {out_dir / 'analysis.json'}")
        print(f"Wrote {out_dir / 'structural_audit.md'}")
        return 0
    except Exception as exc:  # noqa: BLE001 - CLI should show a concise failure.
        print(f"ERROR: {exc}", file=sys.stderr)
        return 2


if __name__ == "__main__":
    raise SystemExit(main())
