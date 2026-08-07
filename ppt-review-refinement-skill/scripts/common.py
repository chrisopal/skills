#!/usr/bin/env python3
"""Shared helpers for PPTX inspection, normalization, and validation."""

from __future__ import annotations

import hashlib
import json
import math
import re
from collections.abc import Iterable, Iterator
from pathlib import Path
from typing import Any

import yaml
from jsonschema import Draft202012Validator
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu
from pptx.oxml.ns import qn

EMU_PER_INCH = 914_400
EMU_PER_POINT = 12_700
HEX_RE = re.compile(r"^#?([0-9A-Fa-f]{6})$")
CJK_RE = re.compile(r"[\u3400-\u4dbf\u4e00-\u9fff\uf900-\ufaff]")


def ensure_dir(path: str | Path) -> Path:
    p = Path(path)
    p.mkdir(parents=True, exist_ok=True)
    return p


def load_json(path: str | Path) -> dict[str, Any]:
    with Path(path).open("r", encoding="utf-8") as f:
        data = json.load(f)
    if not isinstance(data, dict):
        raise ValueError(f"Expected a JSON object in {path}")
    return data


def load_validated_json(
    data_path: str | Path,
    schema_path: str | Path,
    *,
    label: str,
) -> dict[str, Any]:
    """Load a JSON object and enforce the skill's declared JSON Schema contract."""
    data = load_json(data_path)
    return validate_json_data(data, schema_path, label=label)


def validate_json_data(
    data: dict[str, Any],
    schema_path: str | Path,
    *,
    label: str,
) -> dict[str, Any]:
    """Enforce a JSON Schema contract for an already loaded object."""
    if not isinstance(data, dict):
        raise ValueError(f"Expected a JSON object for {label}")
    schema = load_json(schema_path)
    validator = Draft202012Validator(schema)
    errors = sorted(validator.iter_errors(data), key=lambda error: tuple(str(part) for part in error.path))
    if errors:
        details = []
        for error in errors[:8]:
            location = ".".join(str(part) for part in error.path) or "<root>"
            details.append(f"{location}: {error.message}")
        suffix = "" if len(errors) <= 8 else f"; and {len(errors) - 8} more"
        raise ValueError(f"{label} does not match {schema_path}: {'; '.join(details)}{suffix}")
    return data


def write_json(path: str | Path, data: Any) -> None:
    p = Path(path)
    p.parent.mkdir(parents=True, exist_ok=True)
    with p.open("w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
        f.write("\n")


def load_yaml(path: str | Path) -> dict[str, Any]:
    with Path(path).open("r", encoding="utf-8") as f:
        data = yaml.safe_load(f)
    if not isinstance(data, dict):
        raise ValueError(f"Expected a YAML object in {path}")
    return data


def sha256_file(path: str | Path) -> str:
    digest = hashlib.sha256()
    with Path(path).open("rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def sha256_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def emu_to_in(value: int | Emu | None) -> float:
    if value is None:
        return 0.0
    return round(int(value) / EMU_PER_INCH, 4)


def emu_to_pt(value: int | Emu | None) -> float | None:
    if value is None:
        return None
    return round(int(value) / EMU_PER_POINT, 2)


def in_to_emu(value: float) -> int:
    return int(round(value * EMU_PER_INCH))


def normalize_hex(value: str) -> str:
    m = HEX_RE.match(value.strip())
    if not m:
        raise ValueError(f"Invalid RGB hex color: {value!r}")
    return f"#{m.group(1).upper()}"


def hex_to_rgb(value: str) -> RGBColor:
    h = normalize_hex(value)[1:]
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def rgb_to_hex(rgb: RGBColor | bytes | tuple[int, int, int] | None) -> str | None:
    if rgb is None:
        return None
    try:
        vals = list(rgb)
    except TypeError:
        return None
    if len(vals) < 3:
        return None
    return f"#{int(vals[0]):02X}{int(vals[1]):02X}{int(vals[2]):02X}"


def explicit_color_hex(color: ColorFormat | None) -> str | None:
    if color is None:
        return None
    try:
        if color.type == MSO_COLOR_TYPE.RGB:
            return rgb_to_hex(color.rgb)
    except (AttributeError, ValueError, TypeError):
        return None
    return None


def color_descriptor(color: ColorFormat | None) -> str:
    if color is None:
        return "none"
    explicit = explicit_color_hex(color)
    if explicit:
        return explicit
    try:
        if color.type == MSO_COLOR_TYPE.SCHEME:
            return f"theme:{color.theme_color}"
        if color.type is None:
            return "inherit/default"
        return f"type:{color.type}"
    except (AttributeError, ValueError, TypeError):
        return "unknown"


def iter_nested_shapes(shapes: Iterable[Any]) -> Iterator[Any]:
    """Yield shapes recursively. Group child geometry may use a local transform."""
    for shape in shapes:
        yield shape
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from iter_nested_shapes(shape.shapes)
        except (AttributeError, ValueError):
            continue


def iter_text_frames(shape: Any) -> Iterator[Any]:
    """Yield all text frames in a shape, including tables and groups."""
    try:
        if getattr(shape, "has_text_frame", False):
            yield shape.text_frame
    except (AttributeError, ValueError):
        pass

    try:
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell.text_frame
    except (AttributeError, ValueError):
        pass



def iter_runs(shape: Any) -> Iterator[tuple[Any, Any, Any]]:
    """Yield (text_frame, paragraph, run) triples for a shape."""
    for text_frame in iter_text_frames(shape):
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                yield text_frame, paragraph, run


def shape_text(shape: Any) -> str:
    texts: list[str] = []
    for tf in iter_text_frames(shape):
        text = tf.text.strip()
        if text:
            texts.append(text)
    return "\n".join(texts)


def text_overflow_estimate(shape: Any) -> dict[str, float] | None:
    """Estimate obvious text overflow without changing the source presentation.

    This is intentionally conservative: theme-only font sizes and table cells are
    skipped because their effective metrics cannot be inferred reliably here.
    """
    try:
        if not getattr(shape, "has_text_frame", False):
            return None
        text_frame = shape.text_frame
        if not (text_frame.text or "").strip():
            return None
        width_pt = emu_to_pt(shape.width)
        height_pt = emu_to_pt(shape.height)
    except (AttributeError, TypeError, ValueError):
        return None
    if width_pt is None or height_pt is None or width_pt <= 0 or height_pt <= 0:
        return None

    font_sizes = [
        size
        for _text_frame, _paragraph, run in iter_runs(shape)
        if (size := emu_to_pt(run.font.size)) is not None and size > 0
    ]
    if not font_sizes:
        return None
    max_font_pt = max(font_sizes)
    capacity_units = max(1.0, width_pt / max_font_pt)
    estimated_lines = 0
    for paragraph in text_frame.paragraphs:
        paragraph_text = paragraph.text or ""
        if not paragraph_text.strip():
            estimated_lines += 1
            continue
        text_units = sum(
            1.0 if CJK_RE.search(char) else 0.25 if char.isspace() else 0.55
            for char in paragraph_text
        )
        estimated_lines += max(1, math.ceil(text_units / capacity_units))

    line_height_pt = max_font_pt * 1.2
    available_lines = max(0.5, height_pt / line_height_pt)
    ratio = estimated_lines / available_lines
    return {
        "estimated_lines": float(estimated_lines),
        "available_lines": round(available_lines, 2),
        "ratio": round(ratio, 2),
        "font_size_pt": max_font_pt,
    }


def slide_title_shape(slide: Any) -> Any | None:
    try:
        return slide.shapes.title
    except (AttributeError, KeyError, ValueError):
        return None


def text_by_slide(slide: Any) -> dict[str, str]:
    title_shape = slide_title_shape(slide)
    title_id = getattr(title_shape, "shape_id", None)
    title_parts: list[str] = []
    body_parts: list[str] = []

    for shape in iter_nested_shapes(slide.shapes):
        text = shape_text(shape).strip()
        if not text:
            continue
        if getattr(shape, "shape_id", None) == title_id:
            title_parts.append(text)
        else:
            body_parts.append(text)

    return {
        "title": normalize_text("\n".join(title_parts)),
        "body": normalize_text("\n".join(body_parts)),
        "all": normalize_text("\n".join(title_parts + body_parts)),
    }


def normalize_text(text: str) -> str:
    text = text.replace("\u00a0", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\s*\n\s*", "\n", text)
    return text.strip()


def contains_cjk(text: str) -> bool:
    return bool(CJK_RE.search(text))


def set_run_typefaces(run: Any, cjk_font: str, latin_font: str) -> None:
    """Set Latin, East Asian, and complex-script typefaces on a run."""
    r_pr = run._r.get_or_add_rPr()
    for tag, face in (("a:latin", latin_font), ("a:ea", cjk_font), ("a:cs", latin_font)):
        element = r_pr.find(qn(tag))
        if element is None:
            element = OxmlElement(tag)
            r_pr.append(element)
        element.set("typeface", face)
    run.font.name = latin_font


def shape_bbox(shape: Any) -> tuple[float, float, float, float] | None:
    try:
        x = emu_to_in(shape.left)
        y = emu_to_in(shape.top)
        w = emu_to_in(shape.width)
        h = emu_to_in(shape.height)
    except (AttributeError, TypeError, ValueError):
        return None
    if w <= 0 or h <= 0:
        return None
    return x, y, w, h


def bbox_area(box: tuple[float, float, float, float]) -> float:
    return max(0.0, box[2]) * max(0.0, box[3])


def intersection_area(
    a: tuple[float, float, float, float],
    b: tuple[float, float, float, float],
) -> float:
    ax1, ay1, aw, ah = a
    bx1, by1, bw, bh = b
    ax2, ay2 = ax1 + aw, ay1 + ah
    bx2, by2 = bx1 + bw, by1 + bh
    iw = max(0.0, min(ax2, bx2) - max(ax1, bx1))
    ih = max(0.0, min(ay2, by2) - max(ay1, by1))
    return iw * ih


def is_background_like(
    box: tuple[float, float, float, float],
    slide_width_in: float,
    slide_height_in: float,
) -> bool:
    slide_area = slide_width_in * slide_height_in
    return slide_area > 0 and bbox_area(box) / slide_area >= 0.78


def safe_shape_type_name(shape: Any) -> str:
    try:
        return str(shape.shape_type).split(" (")[0]
    except (AttributeError, ValueError):
        return "UNKNOWN"


def safe_shape_name(shape: Any) -> str:
    try:
        return str(shape.name)
    except (AttributeError, ValueError):
        return "unnamed"


def safe_fill_hex(shape: Any) -> str | None:
    try:
        if shape.fill.type == MSO_FILL_TYPE.SOLID:
            return explicit_color_hex(shape.fill.fore_color)
    except (AttributeError, ValueError, TypeError):
        pass
    return None


def safe_line_hex(shape: Any) -> str | None:
    try:
        return explicit_color_hex(shape.line.color)
    except (AttributeError, ValueError, TypeError):
        return None


def luminance(hex_color: str) -> float:
    h = normalize_hex(hex_color)[1:]
    values = [int(h[i : i + 2], 16) / 255.0 for i in (0, 2, 4)]

    def linearize(c: float) -> float:
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4

    r, g, b = (linearize(c) for c in values)
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def contrast_ratio(a: str, b: str) -> float:
    l1, l2 = luminance(a), luminance(b)
    lighter, darker = max(l1, l2), min(l1, l2)
    return round((lighter + 0.05) / (darker + 0.05), 2)


def median(values: list[float]) -> float | None:
    if not values:
        return None
    values = sorted(values)
    n = len(values)
    mid = n // 2
    if n % 2:
        return values[mid]
    return (values[mid - 1] + values[mid]) / 2


def finite(value: float | None) -> bool:
    return value is not None and math.isfinite(value)
