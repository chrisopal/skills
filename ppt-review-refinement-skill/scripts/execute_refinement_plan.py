#!/usr/bin/env python3
"""Execute the explicitly supported, approved subset of an L2/L3 refinement plan."""

from __future__ import annotations

import argparse
import sys
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.util import Pt

from common import (
    explicit_color_hex,
    hex_to_rgb,
    in_to_emu,
    iter_nested_shapes,
    iter_runs,
    load_validated_json,
    safe_shape_name,
    set_run_typefaces,
    sha256_file,
    slide_title_shape,
    write_json,
)
from render_pptx import render_pptx

SKILL_ROOT = Path(__file__).resolve().parents[1]
PLAN_SCHEMA = SKILL_ROOT / "schemas" / "refinement_plan.schema.json"
MANIFEST_SCHEMA = SKILL_ROOT / "schemas" / "change_manifest.schema.json"
TOKENS_SCHEMA = SKILL_ROOT / "schemas" / "style_tokens.schema.json"
PILOT_SCHEMA = SKILL_ROOT / "schemas" / "pilot_confirmation.schema.json"


def _find_shape(slide: Any, shape_id: int | None, *, title: bool = False) -> Any:
    if title:
        shape = slide_title_shape(slide)
        if shape is not None:
            return shape
    if shape_id is None:
        raise ValueError("Executable refinement action requires parameters.shape_id")
    for shape in iter_nested_shapes(slide.shapes):
        if getattr(shape, "shape_id", None) == int(shape_id):
            return shape
    raise ValueError(f"shape_id {shape_id} was not found on slide")


def _same_file_label(declared: str, actual: Path) -> bool:
    declared_path = Path(declared)
    if declared_path.exists():
        return declared_path.resolve() == actual.resolve()
    return declared_path.name == actual.name


def _box(shape: Any) -> dict[str, float]:
    return {
        "x": round(shape.left / 914400, 4),
        "y": round(shape.top / 914400, 4),
        "w": round(shape.width / 914400, 4),
        "h": round(shape.height / 914400, 4),
    }


def _append_change(
    changes: list[dict[str, Any]],
    *,
    slide_num: int,
    shape: Any,
    action: str,
    target: str,
    before: Any,
    after: Any,
    reason: str,
) -> None:
    if before == after:
        return
    changes.append(
        {
            "slide": slide_num,
            "shape_id": getattr(shape, "shape_id", None),
            "shape_name": safe_shape_name(shape),
            "target": target,
            "action": action,
            "before": before,
            "after": after,
            "reason": reason,
        }
    )


def _apply_action(
    slide: Any,
    slide_num: int,
    item: dict[str, Any],
    manifest: dict[str, Any],
    tokens: dict[str, Any],
    changes: list[dict[str, Any]],
) -> None:
    action = str(item.get("action", "")).strip()
    params = item.get("parameters", {})
    if not isinstance(params, dict):
        raise ValueError(f"slide {slide_num} action {action} parameters must be an object")
    target_name = str(item.get("target", ""))
    reason = str(item.get("reason", ""))
    authorization = str(item.get("authorization", "")).strip()
    if not authorization:
        raise PermissionError(f"slide {slide_num} action {action} has no authorization reference")
    visual_policy = manifest.get("visual_policy", {})
    policy_by_action = {
        "set_geometry": ("layout", {"optimize", "restructure"}),
        "set_title_geometry": ("layout", {"optimize", "restructure"}),
        "set_font_role": ("typography", {"normalize", "redesign"}),
        "set_fill_color": ("color", {"normalize", "redesign"}),
        "set_line_color": ("color", {"normalize", "redesign"}),
    }
    if action in policy_by_action:
        policy_name, allowed_values = policy_by_action[action]
        if visual_policy.get(policy_name) not in allowed_values:
            raise PermissionError(
                f"slide {slide_num} action {action} is outside manifest visual_policy.{policy_name}"
            )
    title_target = target_name.lower() in {"title", "title_box", "title geometry"}
    shape = _find_shape(slide, params.get("shape_id"), title=title_target)

    if action in {"set_geometry", "set_title_geometry"}:
        if action == "set_title_geometry" and not title_target:
            shape = _find_shape(slide, params.get("shape_id"), title=True)
        required = {"x", "y", "w", "h"}
        missing = required - set(params)
        if missing:
            raise ValueError(f"slide {slide_num} {action} missing parameters: {sorted(missing)}")
        before = _box(shape)
        for key in required:
            if float(params[key]) < 0:
                raise ValueError(f"slide {slide_num} {action} does not accept negative geometry")
        shape.left = in_to_emu(float(params["x"]))
        shape.top = in_to_emu(float(params["y"]))
        shape.width = in_to_emu(float(params["w"]))
        shape.height = in_to_emu(float(params["h"]))
        _append_change(changes, slide_num=slide_num, shape=shape, action=action, target=target_name, before=before, after=_box(shape), reason=reason)
        return

    if action == "set_font_role":
        role_name = str(params.get("role", ""))
        role = tokens.get("typography", {}).get("roles", {}).get(role_name)
        if not isinstance(role, dict):
            raise ValueError(f"Unknown typography role: {role_name}")
        fonts = tokens.get("typography", {}).get("fonts", {})
        cjk_font = str(fonts.get("cjk", "Microsoft YaHei"))
        latin_font = str(fonts.get("latin", cjk_font))
        for _text_frame, paragraph, run in iter_runs(shape):
            if not (run.text or "").strip():
                continue
            before = {
                "size_pt": run.font.size.pt if run.font.size is not None else None,
                "bold": run.font.bold,
                "color": explicit_color_hex(run.font.color),
            }
            set_run_typefaces(run, cjk_font, latin_font)
            if role.get("size_pt") is not None:
                run.font.size = Pt(float(role["size_pt"]))
            if role.get("bold") is not None:
                run.font.bold = bool(role["bold"])
            if role.get("color"):
                run.font.color.rgb = hex_to_rgb(str(role["color"]))
            if role.get("line_spacing") is not None:
                paragraph.line_spacing = float(role["line_spacing"])
            if role.get("space_after_pt") is not None:
                paragraph.space_after = Pt(float(role["space_after_pt"]))
            after = {
                "size_pt": run.font.size.pt if run.font.size is not None else None,
                "bold": run.font.bold,
                "color": explicit_color_hex(run.font.color),
            }
            _append_change(changes, slide_num=slide_num, shape=shape, action=action, target=target_name, before=before, after=after, reason=reason)
        return

    if action == "set_fill_color":
        color = str(params.get("color", ""))
        before = explicit_color_hex(shape.fill.fore_color) if shape.fill.type == MSO_FILL_TYPE.SOLID else None
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(color)
        _append_change(changes, slide_num=slide_num, shape=shape, action=action, target=target_name, before=before, after=color.upper(), reason=reason)
        return

    if action == "set_line_color":
        color = str(params.get("color", ""))
        before = explicit_color_hex(shape.line.color)
        shape.line.color.rgb = hex_to_rgb(color)
        _append_change(changes, slide_num=slide_num, shape=shape, action=action, target=target_name, before=before, after=color.upper(), reason=reason)
        return

    if action == "set_text":
        if not getattr(shape, "has_text_frame", False):
            raise ValueError(f"slide {slide_num} set_text target is not a text shape")
        if not title_target or manifest.get("content_policy", {}).get("titles") not in {"optimize", "rewrite"}:
            raise PermissionError("set_text is limited to explicitly authorized title changes")
        text = str(params.get("text", ""))
        before = shape.text
        shape.text_frame.clear()
        shape.text_frame.paragraphs[0].add_run().text = text
        _append_change(changes, slide_num=slide_num, shape=shape, action=action, target=target_name, before=before, after=text, reason=reason)
        return

    raise ValueError(f"Unsupported refinement action on slide {slide_num}: {action}")


def execute_refinement_plan(
    input_path: str | Path,
    output_path: str | Path,
    plan_path: str | Path,
    tokens_path: str | Path,
    manifest_path: str | Path,
    pilot_confirmation_path: str | Path | None,
    *,
    log_path: str | Path | None = None,
    before_after_dir: str | Path | None = None,
    render: bool = True,
) -> dict[str, Any]:
    input_path = Path(input_path).resolve()
    output_path = Path(output_path).resolve()
    if input_path == output_path:
        raise ValueError("Refusing to overwrite the source deck")
    if input_path.suffix.lower() != ".pptx" or output_path.suffix.lower() != ".pptx":
        raise ValueError("L2/L3 executor supports .pptx only; .pptm needs an Office-level executor")
    plan = load_validated_json(plan_path, PLAN_SCHEMA, label="refinement_plan.json")
    tokens = load_validated_json(tokens_path, TOKENS_SCHEMA, label="style_tokens.json")
    manifest = load_validated_json(manifest_path, MANIFEST_SCHEMA, label="change_manifest.json")
    level = str(plan["optimization_level"])
    declared_source = Path(str(plan["source_file"]))
    if declared_source.exists():
        source_matches = declared_source.resolve() == input_path
    else:
        source_matches = declared_source.name == input_path.name
    if not source_matches:
        raise ValueError("refinement_plan.source_file does not match the input deck")
    if level not in {"L2", "L3"}:
        raise PermissionError("L2/L3 executor only accepts L2 or L3 refinement plans")
    if manifest.get("approval_status") != "approved":
        raise PermissionError("change_manifest.json must be approved before L2/L3 execution")
    if manifest.get("optimization_level") != level:
        raise ValueError("refinement_plan and change_manifest optimization levels must match")
    if pilot_confirmation_path is None:
        raise PermissionError("L2/L3 execution requires an approved pilot_confirmation.json")
    pilot = load_validated_json(pilot_confirmation_path, PILOT_SCHEMA, label="pilot_confirmation.json")
    if pilot.get("status") != "approved":
        raise PermissionError("pilot_confirmation.json status is not approved")
    if not _same_file_label(pilot["source_file"], input_path):
        raise ValueError("pilot_confirmation.json source_file does not match the input deck")
    if pilot["source_sha256"].lower() != sha256_file(input_path):
        raise ValueError("pilot_confirmation.json source_sha256 does not match the input deck")
    plan_slides = {int(item["slide"]) for item in plan["slides"]}
    if not plan_slides.issubset(set(int(value) for value in pilot["reviewed_slides"])):
        raise PermissionError("pilot_confirmation.json does not cover every planned slide")
    prs = Presentation(str(input_path))
    if any(slide_num > len(prs.slides) for slide_num in plan_slides):
        raise ValueError("refinement_plan references a slide outside the source deck")
    changes: list[dict[str, Any]] = []
    for slide_item in plan["slides"]:
        slide_num = int(slide_item["slide"])
        if slide_num in set(int(value) for value in manifest.get("protected_slides", [])) and (
            slide_item.get("change") or slide_item.get("content_changes") or slide_item.get("image_actions")
        ):
            raise PermissionError(f"refinement_plan attempts to change protected slide {slide_num}")
        slide = prs.slides[slide_num - 1]
        for item in slide_item.get("change", []) + slide_item.get("content_changes", []):
            _apply_action(slide, slide_num, item, manifest, tokens, changes)
        if slide_item.get("image_actions"):
            raise ValueError(f"slide {slide_num} contains image_actions; image replacement is outside this executor")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    Presentation(str(output_path))
    rendered = {}
    if render:
        render_root = Path(before_after_dir or output_path.parent / f"{output_path.stem}-before-after")
        before_dir = render_root / "before"
        after_dir = render_root / "after"
        rendered["before"] = render_pptx(input_path, before_dir)
        rendered["after"] = render_pptx(output_path, after_dir)

    log = {
        "version": "1.0.0",
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "source_file": str(input_path),
        "source_sha256": sha256_file(input_path),
        "output_file": str(output_path),
        "output_sha256": sha256_file(output_path),
        "plan_file": str(Path(plan_path).resolve()),
        "tokens_file": str(Path(tokens_path).resolve()),
        "manifest_file": str(Path(manifest_path).resolve()),
        "pilot_confirmation_file": str(Path(pilot_confirmation_path).resolve()),
        "optimization_level": level,
        "change_count": len(changes),
        "changes": changes,
        "rendered_before_after": rendered,
        "limitations": [
            "Only set_geometry, set_title_geometry, set_font_role, set_fill_color, set_line_color, and authorized title set_text are executable.",
            "SmartArt, charts, animations, embedded objects, image replacement, and macro-enabled decks remain outside this executor.",
            "Final visual quality still requires rendered review and visual_signoff.json.",
        ],
    }
    if log_path:
        write_json(log_path, log)
    return log


def main() -> int:
    parser = argparse.ArgumentParser(description="Execute an approved L2/L3 refinement plan.")
    parser.add_argument("--input", required=True)
    parser.add_argument("--output", required=True)
    parser.add_argument("--plan", required=True)
    parser.add_argument("--tokens", required=True)
    parser.add_argument("--manifest", required=True)
    parser.add_argument("--pilot-confirmation", required=True)
    parser.add_argument("--log", required=True)
    parser.add_argument("--before-after-dir")
    parser.add_argument("--no-render", action="store_true")
    args = parser.parse_args()
    try:
        log = execute_refinement_plan(
            args.input,
            args.output,
            args.plan,
            args.tokens,
            args.manifest,
            args.pilot_confirmation,
            log_path=args.log,
            before_after_dir=args.before_after_dir,
            render=not args.no_render,
        )
        print(f"Wrote {args.output} with {log['change_count']} logged changes")
        return 0
    except Exception as exc:  # noqa: BLE001
        print(f"ERROR: {exc}", file=sys.stderr)
        return 2


if __name__ == "__main__":
    raise SystemExit(main())
