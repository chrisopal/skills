#!/usr/bin/env python3
"""Apply low-risk, explicitly authorized normalization to a PPTX deck."""

from __future__ import annotations

import argparse
import sys
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.util import Pt

from common import (
    explicit_color_hex,
    hex_to_rgb,
    in_to_emu,
    iter_nested_shapes,
    iter_runs,
    load_validated_json,
    normalize_hex,
    safe_shape_name,
    sha256_file,
    slide_title_shape,
    set_run_typefaces,
    write_json,
)

SKILL_ROOT = Path(__file__).resolve().parents[1]
MANIFEST_SCHEMA = SKILL_ROOT / "schemas" / "change_manifest.schema.json"
TOKENS_SCHEMA = SKILL_ROOT / "schemas" / "style_tokens.schema.json"


def normalize_mapping(mapping: dict[str, Any]) -> dict[str, str]:
    result: dict[str, str] = {}
    for source, target in mapping.items():
        result[normalize_hex(str(source))] = normalize_hex(str(target))
    return result


def add_change(
    changes: list[dict[str, Any]],
    *,
    slide: int,
    shape: Any,
    property_name: str,
    before: Any,
    after: Any,
    reason: str,
    text_preview: str = "",
) -> None:
    if before == after:
        return
    changes.append(
        {
            "slide": slide,
            "shape_id": getattr(shape, "shape_id", None),
            "shape_name": safe_shape_name(shape),
            "property": property_name,
            "before": before,
            "after": after,
            "reason": reason,
            "text_preview": text_preview[:100],
        }
    )


def apply_fill_mapping(
    fill: Any,
    mapping: dict[str, str],
    *,
    slide_num: int,
    shape: Any,
    property_name: str,
    changes: list[dict[str, Any]],
) -> None:
    try:
        if fill.type != MSO_FILL_TYPE.SOLID:
            return
        current = explicit_color_hex(fill.fore_color)
        if not current or current not in mapping:
            return
        target = mapping[current]
        fill.fore_color.rgb = hex_to_rgb(target)
        add_change(
            changes,
            slide=slide_num,
            shape=shape,
            property_name=property_name,
            before=current,
            after=target,
            reason="exact_color_mapping",
        )
    except (AttributeError, ValueError, TypeError):
        return


def normalize_pptx(
    input_path: str | Path,
    output_path: str | Path,
    tokens_path: str | Path,
    manifest_path: str | Path,
    *,
    allow_draft: bool = False,
    apply_role_sizes: bool = False,
    standardize_title_position: bool = False,
) -> dict[str, Any]:
    input_path = Path(input_path).resolve()
    output_path = Path(output_path).resolve()
    if input_path == output_path:
        raise ValueError("Refusing to overwrite the source deck. Choose a different output path.")
    if not input_path.exists():
        raise FileNotFoundError(input_path)
    if input_path.suffix.lower() == ".pptm" or output_path.suffix.lower() == ".pptm":
        raise ValueError(
            "Safe normalizer refuses .pptm because python-pptx cannot guarantee VBA preservation. "
            "Use a PowerPoint/Office API executor for macro-enabled decks."
        )
    if input_path.suffix.lower() != ".pptx" or output_path.suffix.lower() != ".pptx":
        raise ValueError("Safe normalization supports .pptx output only.")
    output_path.parent.mkdir(parents=True, exist_ok=True)

    tokens = load_validated_json(tokens_path, TOKENS_SCHEMA, label="style_tokens.json")
    manifest = load_validated_json(manifest_path, MANIFEST_SCHEMA, label="change_manifest.json")
    approval_status = manifest.get("approval_status", "draft")
    if approval_status != "approved" and not allow_draft:
        raise PermissionError(
            "change_manifest.json is not approved. Set approval_status to 'approved' or use --allow-draft for testing."
        )
    if manifest.get("optimization_level") == "L0":
        raise PermissionError("L0 is review-only; normalization is not authorized.")

    visual_policy = manifest.get("visual_policy", {})
    typography_allowed = visual_policy.get("typography") in {"normalize", "redesign"}
    color_allowed = visual_policy.get("color") in {"normalize", "redesign"}
    layout_allowed = visual_policy.get("layout") in {"optimize", "restructure"}
    protected_slides = set(int(v) for v in manifest.get("protected_slides", []))

    fonts = tokens.get("typography", {}).get("fonts", {})
    cjk_font = str(fonts.get("cjk", "Microsoft YaHei"))
    latin_font = str(fonts.get("latin", cjk_font))
    roles = tokens.get("typography", {}).get("roles", {})
    exact_mappings = tokens.get("exact_mappings", {})
    font_color_map = normalize_mapping(exact_mappings.get("font_colors", {}))
    fill_color_map = normalize_mapping(exact_mappings.get("fill_colors", {}))
    line_color_map = normalize_mapping(exact_mappings.get("line_colors", {}))

    prs = Presentation(str(input_path))
    changes: list[dict[str, Any]] = []
    skipped: list[dict[str, Any]] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        if slide_num in protected_slides:
            skipped.append({"slide": slide_num, "reason": "protected_slide"})
            continue

        title_shape = slide_title_shape(slide)
        title_shape_id = getattr(title_shape, "shape_id", None)

        for shape in iter_nested_shapes(slide.shapes):
            is_title = getattr(shape, "shape_id", None) == title_shape_id
            role_name = "cover_title" if slide_num == 1 and is_title else "page_title" if is_title else "body"
            role = roles.get(role_name, {})

            if typography_allowed:
                for _tf, _paragraph, run in iter_runs(shape):
                    if not (run.text or "").strip():
                        continue
                    before_font = run.font.name or "<theme/default>"
                    set_run_typefaces(run, cjk_font, latin_font)
                    add_change(
                        changes,
                        slide=slide_num,
                        shape=shape,
                        property_name="font_family",
                        before=before_font,
                        after=f"latin={latin_font}; eastAsian={cjk_font}",
                        reason="typography_normalization",
                        text_preview=run.text,
                    )

                    if apply_role_sizes and is_title and role:
                        target_size = role.get("size_pt")
                        target_bold = role.get("bold")
                        if target_size is not None:
                            before_size = run.font.size.pt if run.font.size is not None else None
                            run.font.size = Pt(float(target_size))
                            add_change(
                                changes,
                                slide=slide_num,
                                shape=shape,
                                property_name="font_size_pt",
                                before=before_size,
                                after=float(target_size),
                                reason=f"role_style:{role_name}",
                                text_preview=run.text,
                            )
                        if target_bold is not None:
                            before_bold = run.font.bold
                            run.font.bold = bool(target_bold)
                            add_change(
                                changes,
                                slide=slide_num,
                                shape=shape,
                                property_name="font_bold",
                                before=before_bold,
                                after=bool(target_bold),
                                reason=f"role_style:{role_name}",
                                text_preview=run.text,
                            )

                    if color_allowed:
                        current_font_color = explicit_color_hex(run.font.color)
                        if current_font_color and current_font_color in font_color_map:
                            target = font_color_map[current_font_color]
                            run.font.color.rgb = hex_to_rgb(target)
                            add_change(
                                changes,
                                slide=slide_num,
                                shape=shape,
                                property_name="font_color",
                                before=current_font_color,
                                after=target,
                                reason="exact_color_mapping",
                                text_preview=run.text,
                            )

            if color_allowed:
                try:
                    shape_fill = shape.fill
                except (AttributeError, ValueError, TypeError):
                    shape_fill = None
                if shape_fill is not None:
                    apply_fill_mapping(
                        shape_fill,
                        fill_color_map,
                        slide_num=slide_num,
                        shape=shape,
                        property_name="fill_color",
                        changes=changes,
                    )
                try:
                    current_line_color = explicit_color_hex(shape.line.color)
                    if current_line_color and current_line_color in line_color_map:
                        target = line_color_map[current_line_color]
                        shape.line.color.rgb = hex_to_rgb(target)
                        add_change(
                            changes,
                            slide=slide_num,
                            shape=shape,
                            property_name="line_color",
                            before=current_line_color,
                            after=target,
                            reason="exact_color_mapping",
                        )
                except (AttributeError, ValueError, TypeError):
                    pass

                try:
                    if getattr(shape, "has_table", False):
                        for row_idx, row in enumerate(shape.table.rows):
                            for col_idx, cell in enumerate(row.cells):
                                before = None
                                try:
                                    if cell.fill.type == MSO_FILL_TYPE.SOLID:
                                        before = explicit_color_hex(cell.fill.fore_color)
                                except (AttributeError, ValueError, TypeError):
                                    before = None
                                if before and before in fill_color_map:
                                    target = fill_color_map[before]
                                    cell.fill.fore_color.rgb = hex_to_rgb(target)
                                    add_change(
                                        changes,
                                        slide=slide_num,
                                        shape=shape,
                                        property_name=f"table_cell_fill[{row_idx},{col_idx}]",
                                        before=before,
                                        after=target,
                                        reason="exact_color_mapping",
                                    )
                except (AttributeError, ValueError, TypeError):
                    pass

        if standardize_title_position and layout_allowed and title_shape is not None and slide_num > 1:
            box = tokens.get("layout", {}).get("title_box_in", {})
            if all(key in box for key in ("x", "y", "w", "h")):
                before_box = {
                    "x": round(title_shape.left / 914400, 4),
                    "y": round(title_shape.top / 914400, 4),
                    "w": round(title_shape.width / 914400, 4),
                    "h": round(title_shape.height / 914400, 4),
                }
                title_shape.left = in_to_emu(float(box["x"]))
                title_shape.top = in_to_emu(float(box["y"]))
                title_shape.width = in_to_emu(float(box["w"]))
                title_shape.height = in_to_emu(float(box["h"]))
                after_box = {key: float(box[key]) for key in ("x", "y", "w", "h")}
                add_change(
                    changes,
                    slide=slide_num,
                    shape=title_shape,
                    property_name="title_box_in",
                    before=before_box,
                    after=after_box,
                    reason="layout_token:title_box_in",
                    text_preview=getattr(title_shape, "text", ""),
                )

    prs.save(str(output_path))
    # Re-open to catch corrupt serialization before returning.
    Presentation(str(output_path))

    log = {
        "version": "1.0.0",
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "source_file": str(input_path),
        "source_sha256": sha256_file(input_path),
        "output_file": str(output_path),
        "output_sha256": sha256_file(output_path),
        "tokens_file": str(Path(tokens_path).resolve()),
        "manifest_file": str(Path(manifest_path).resolve()),
        "approval_status": approval_status,
        "options": {
            "apply_role_sizes": apply_role_sizes,
            "standardize_title_position": standardize_title_position,
        },
        "change_count": len(changes),
        "changes": changes,
        "skipped": skipped,
        "limitations": [
            "This safe normalizer does not rewrite text, reorder slides, replace images, restyle charts, or rebuild diagrams.",
            "Theme-inherited colors cannot be remapped unless they are explicit RGB values.",
            "Animations, SmartArt, and embedded objects require PowerPoint-level manual verification.",
        ],
    }
    return log


def main() -> int:
    parser = argparse.ArgumentParser(description="Apply safe, manifest-controlled PPTX normalization.")
    parser.add_argument("input", help="Source PPTX/PPTM")
    parser.add_argument("output", help="Output PPTX/PPTM; must differ from source")
    parser.add_argument("--tokens", required=True, help="style_tokens.json")
    parser.add_argument("--manifest", required=True, help="approved change_manifest.json")
    parser.add_argument("--log", required=True, help="Output change log JSON")
    parser.add_argument("--allow-draft", action="store_true", help="Allow draft manifest for testing only")
    parser.add_argument("--apply-role-sizes", action="store_true", help="Apply title role size/bold")
    parser.add_argument("--standardize-title-position", action="store_true", help="Apply title box token to content slides")
    args = parser.parse_args()

    try:
        log = normalize_pptx(
            args.input,
            args.output,
            args.tokens,
            args.manifest,
            allow_draft=args.allow_draft,
            apply_role_sizes=args.apply_role_sizes,
            standardize_title_position=args.standardize_title_position,
        )
        write_json(args.log, log)
        print(f"Wrote {args.output} with {log['change_count']} logged changes")
        print(f"Wrote {args.log}")
        return 0
    except Exception as exc:  # noqa: BLE001
        print(f"ERROR: {exc}", file=sys.stderr)
        return 2


if __name__ == "__main__":
    raise SystemExit(main())
