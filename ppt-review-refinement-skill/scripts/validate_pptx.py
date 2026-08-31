#!/usr/bin/env python3
"""Validate an optimized PPTX against its source and approved change manifest."""

from __future__ import annotations

import argparse
import re
import sys
from collections import Counter
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from analyze_pptx import analyze_deck
from confirm_visual_review import validate_visual_signoff
from common import (
    ensure_dir,
    iter_nested_shapes,
    load_validated_json,
    normalize_text,
    sha256_bytes,
    sha256_file,
    text_by_slide,
    write_json,
)
from make_montage import create_montage
from render_pptx import render_pptx

NUMBER_RE = re.compile(
    r"(?<![A-Za-z0-9_])(?:[¥￥$€£]?\s*)?\d+(?:[,.]\d+)*(?:\s*(?:%|％|亿元|万元|亿|万|元|美元|欧元|"
    r"年|月|日|天|小时|分钟|秒|件|台|套|条|人|家|个|次|页|章|GB|TB|MB|KB|GHz|MHz|kW|MW|kWh|"
    r"kg|g|t|mm|cm|m|km|㎡|m²|°C|℃|pt|DPI))?(?![A-Za-z0-9_])",
    flags=re.IGNORECASE,
)

SKILL_ROOT = Path(__file__).resolve().parents[1]
MANIFEST_SCHEMA = SKILL_ROOT / "schemas" / "change_manifest.schema.json"


def make_check(check_id: str, name: str, status: str, details: str) -> dict[str, str]:
    return {"id": check_id, "name": name, "status": status, "details": details}


def numeric_tokens(text: str) -> Counter[str]:
    tokens: list[str] = []
    for match in NUMBER_RE.finditer(text):
        token = re.sub(r"\s+", "", match.group(0))
        token = token.replace("，", ",").replace("％", "%")
        tokens.append(token)
    return Counter(tokens)


def image_hashes_by_slide(prs: Presentation) -> list[Counter[str]]:
    result: list[Counter[str]] = []
    for slide in prs.slides:
        hashes: Counter[str] = Counter()
        for shape in iter_nested_shapes(slide.shapes):
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    hashes[sha256_bytes(shape.image.blob)] += 1
            except (AttributeError, ValueError, KeyError):
                continue
        result.append(hashes)
    return result


def counter_diff(before: Counter[str], after: Counter[str]) -> dict[str, dict[str, int]]:
    keys = sorted(set(before) | set(after))
    return {
        key: {"before": before.get(key, 0), "after": after.get(key, 0)}
        for key in keys
        if before.get(key, 0) != after.get(key, 0)
    }


def accepted_risk_ids(manifest: dict[str, Any]) -> set[str]:
    return {str(item.get("id")) for item in manifest.get("accepted_risks", []) if item.get("id")}


def validate_deck(
    original_path: str | Path,
    candidate_path: str | Path,
    manifest_path: str | Path,
    output_dir: str | Path,
    *,
    render: bool = True,
    visual_signoff_path: str | Path | None = None,
    allow_unconfirmed_visual: bool = False,
) -> dict[str, Any]:
    original_path = Path(original_path).resolve()
    candidate_path = Path(candidate_path).resolve()
    output_dir = ensure_dir(output_dir).resolve()
    manifest = load_validated_json(manifest_path, MANIFEST_SCHEMA, label="change_manifest.json")
    accepted = accepted_risk_ids(manifest)
    checks: list[dict[str, str]] = []
    differences: list[dict[str, Any]] = []
    issues: list[dict[str, Any]] = []
    manual_review: list[str] = []

    try:
        original = Presentation(str(original_path))
        candidate = Presentation(str(candidate_path))
        checks.append(make_check("VAL-OPEN", "文件可打开", "pass", "原始文件和候选文件均可解析。"))
    except Exception as exc:  # noqa: BLE001
        checks.append(make_check("VAL-OPEN", "文件可打开", "fail", str(exc)))
        report = {
            "version": "1.0.0",
            "generated_at": datetime.now(timezone.utc).isoformat(),
            "status": "fail",
            "checks": checks,
            "differences": differences,
            "issues": issues,
            "manual_review_required": manual_review,
            "metrics": {},
        }
        return report

    original_count = len(original.slides)
    candidate_count = len(candidate.slides)
    if visual_signoff_path:
        try:
            signoff = validate_visual_signoff(visual_signoff_path, candidate_path, original_path)
            checks.append(
                make_check(
                    "VAL-VISUAL-SIGNOFF",
                    "最终人工视觉确认",
                    "pass",
                    f"已由 {signoff['reviewer']} 确认全部 {len(signoff['reviewed_slides'])} 页及六项视觉检查。",
                )
            )
        except Exception as exc:  # noqa: BLE001
            checks.append(make_check("VAL-VISUAL-SIGNOFF", "最终人工视觉确认", "fail", str(exc)))
            manual_review.append("最终人工视觉确认文件无效或未覆盖候选文件全部页面。")
    elif allow_unconfirmed_visual:
        checks.append(
            make_check(
                "VAL-VISUAL-SIGNOFF",
                "最终人工视觉确认",
                "pass",
                "按显式测试选项跳过最终人工视觉确认；此结果不可作为最终交付证明。",
            )
        )
        manual_review.append("最终人工视觉确认尚未完成。")
    else:
        checks.append(
            make_check(
                "VAL-VISUAL-SIGNOFF",
                "最终人工视觉确认",
                "fail",
                "未提供 approved visual_signoff.json；最终交付必须包含人工视觉确认。",
            )
        )
        manual_review.append("必须由人工查看渲染结果并提交 approved visual_signoff.json。")
    slide_order_policy = manifest.get("content_policy", {}).get("slide_order", "preserve")
    split_policy = manifest.get("content_policy", {}).get("split_merge_delete", "forbidden")
    if original_count == candidate_count:
        checks.append(make_check("VAL-SLIDES", "页数", "pass", f"页数保持为 {candidate_count}。"))
    elif split_policy == "allowed":
        checks.append(
            make_check(
                "VAL-SLIDES",
                "页数",
                "warning",
                f"页数由 {original_count} 变为 {candidate_count}，授权允许拆分/合并，但需要人工确认顺序和内容。",
            )
        )
        differences.append({"type": "slide_count", "before": original_count, "after": candidate_count})
        manual_review.append("页面拆分/合并后的叙事顺序和内容完整性需要人工对照。")
    else:
        checks.append(
            make_check(
                "VAL-SLIDES",
                "页数",
                "fail",
                f"页数由 {original_count} 变为 {candidate_count}，未获得拆分/合并授权。",
            )
        )
        differences.append({"type": "slide_count", "before": original_count, "after": candidate_count})

    original_text = [text_by_slide(slide) for slide in original.slides]
    candidate_text = [text_by_slide(slide) for slide in candidate.slides]
    compare_count = min(original_count, candidate_count)
    content_policy = manifest.get("content_policy", {})
    body_policy = content_policy.get("body_text", "preserve")
    title_policy = content_policy.get("titles", "preserve")
    protected_slides = set(int(v) for v in manifest.get("protected_slides", []))

    body_failures = 0
    title_failures = 0
    protected_slide_failures = 0
    for index in range(compare_count):
        slide_num = index + 1
        before = original_text[index]
        after = candidate_text[index]
        if slide_num in protected_slides and before["all"] != after["all"]:
            protected_slide_failures += 1
            differences.append(
                {
                    "type": "protected_slide_text",
                    "slide": slide_num,
                    "before": before["all"],
                    "after": after["all"],
                }
            )
        if body_policy == "preserve" and before["body"] != after["body"]:
            body_failures += 1
            differences.append(
                {"type": "body_text", "slide": slide_num, "before": before["body"], "after": after["body"]}
            )
        if title_policy == "preserve" and before["title"] != after["title"]:
            title_failures += 1
            differences.append(
                {"type": "title_text", "slide": slide_num, "before": before["title"], "after": after["title"]}
            )

    checks.append(
        make_check(
            "VAL-BODY",
            "正文保护",
            "pass" if body_failures == 0 else "fail",
            "正文文本未变化。" if body_failures == 0 else f"有 {body_failures} 页正文发生未授权变化。",
        )
    )
    checks.append(
        make_check(
            "VAL-TITLE",
            "标题保护",
            "pass" if title_failures == 0 else "fail",
            (
                "标题按授权规则通过。"
                if title_failures == 0
                else f"有 {title_failures} 页标题发生未授权变化。"
            ),
        )
    )
    checks.append(
        make_check(
            "VAL-PROTECTED-SLIDES",
            "受保护页面",
            "pass" if protected_slide_failures == 0 else "fail",
            (
                "受保护页面文本未变化。"
                if protected_slide_failures == 0
                else f"有 {protected_slide_failures} 个受保护页面发生变化。"
            ),
        )
    )

    numeric_failures = 0
    if slide_order_policy == "preserve" and original_count == candidate_count:
        for index in range(compare_count):
            before_tokens = numeric_tokens(original_text[index]["all"])
            after_tokens = numeric_tokens(candidate_text[index]["all"])
            diff = counter_diff(before_tokens, after_tokens)
            if diff:
                numeric_failures += 1
                differences.append({"type": "numeric_tokens", "slide": index + 1, "differences": diff})
    else:
        before_tokens = numeric_tokens("\n".join(item["all"] for item in original_text))
        after_tokens = numeric_tokens("\n".join(item["all"] for item in candidate_text))
        diff = counter_diff(before_tokens, after_tokens)
        if diff:
            numeric_failures = 1
            differences.append({"type": "numeric_tokens_global", "differences": diff})

    checks.append(
        make_check(
            "VAL-NUMBERS",
            "数字、日期与单位保护",
            "pass" if numeric_failures == 0 else "fail",
            "数字 Token 未变化。" if numeric_failures == 0 else f"检测到 {numeric_failures} 处数字 Token 差异。",
        )
    )

    protected_term_failures = 0
    original_all = "\n".join(item["all"] for item in original_text)
    candidate_all = "\n".join(item["all"] for item in candidate_text)
    for term in manifest.get("protected_terms", []):
        before_count = original_all.count(term)
        after_count = candidate_all.count(term)
        if before_count != after_count:
            protected_term_failures += 1
            differences.append(
                {
                    "type": "protected_term",
                    "term": term,
                    "before_count": before_count,
                    "after_count": after_count,
                }
            )
    checks.append(
        make_check(
            "VAL-TERMS",
            "受保护术语",
            "pass" if protected_term_failures == 0 else "fail",
            (
                "受保护术语出现次数未变化。"
                if protected_term_failures == 0
                else f"有 {protected_term_failures} 个受保护术语发生变化。"
            ),
        )
    )

    image_policy = manifest.get("image_policy", {}).get("replace", "approval_required")
    image_diffs = 0
    if original_count == candidate_count:
        before_images = image_hashes_by_slide(original)
        after_images = image_hashes_by_slide(candidate)
        for index in range(compare_count):
            diff = counter_diff(before_images[index], after_images[index])
            if diff:
                image_diffs += 1
                differences.append({"type": "image_hashes", "slide": index + 1, "differences": diff})
    if image_diffs == 0:
        checks.append(
            make_check(
                "VAL-IMAGES",
                "图片替换保护",
                "pass",
                "图片二进制内容未变化。",
            )
        )
    elif image_policy == "forbidden":
        checks.append(
            make_check(
                "VAL-IMAGES",
                "图片替换保护",
                "fail",
                f"有 {image_diffs} 页图片发生变化，但 manifest 禁止替换图片。",
            )
        )
    else:
        image_check_status = "pass" if "VAL-IMAGES" in accepted else "warning"
        checks.append(
            make_check(
                "VAL-IMAGES",
                "图片替换保护",
                image_check_status,
                (
                    "图片替换风险已在 accepted_risks 中登记；仍需根据变更日志人工核对来源和真实性。"
                    if image_check_status == "pass"
                    else "授权允许或要求逐页确认图片替换；需要在 accepted_risks 中登记后才能交付。"
                ),
            )
        )
        manual_review.append("逐页核对被替换或重新生成的图片、来源、授权和真实性标识。")

    candidate_analysis = analyze_deck(candidate_path)
    p0_issues = [item for item in candidate_analysis["issues"] if item["severity"] == "P0"]
    p1_issues = [item for item in candidate_analysis["issues"] if item["severity"] == "P1"]
    issues.extend(candidate_analysis["issues"])
    if p0_issues:
        checks.append(
            make_check(
                "VAL-STRUCTURE",
                "结构与几何风险",
                "fail",
                f"候选文件仍有 {len(p0_issues)} 个 P0 和 {len(p1_issues)} 个 P1 结构风险。",
            )
        )
    elif p1_issues:
        risk_id = "VAL-STRUCTURE-P1"
        status = "warning" if risk_id in accepted else "fail"
        checks.append(
            make_check(
                "VAL-STRUCTURE",
                "结构与几何风险",
                status,
                f"候选文件有 {len(p1_issues)} 个 P1 风险。要以风险方式交付，需在 accepted_risks 中加入 {risk_id}。",
            )
        )
    else:
        checks.append(make_check("VAL-STRUCTURE", "结构与几何风险", "pass", "未检测到 P0/P1 结构风险。"))

    render_status = "not_checked"
    rendered_count = 0
    if render:
        try:
            render_dir = ensure_dir(output_dir / "rendered")
            render_result = render_pptx(candidate_path, render_dir)
            rendered_count = int(render_result["slide_count"])
            create_montage(render_dir, output_dir / "candidate_montage.png")
            if rendered_count == candidate_count:
                render_status = "pass"
                checks.append(
                    make_check(
                        "VAL-RENDER",
                        "页面渲染",
                        "pass",
                        f"候选文件成功渲染 {rendered_count} 页。仍需人工逐页查看 PNG。",
                    )
                )
                manual_review.append("逐页查看 rendered/slide-*.png，确认无视觉溢出、错位、变形和字体替换。")
            else:
                render_status = "fail"
                checks.append(
                    make_check(
                        "VAL-RENDER",
                        "页面渲染",
                        "fail",
                        f"候选文件共 {candidate_count} 页，但只渲染出 {rendered_count} 页。",
                    )
                )
        except Exception as exc:  # noqa: BLE001
            render_status = "fail"
            checks.append(make_check("VAL-RENDER", "页面渲染", "fail", str(exc)))
    else:
        status = "warning" if "VAL-RENDER-SKIPPED" in accepted else "fail"
        checks.append(
            make_check(
                "VAL-RENDER-SKIPPED",
                "页面渲染",
                status,
                "本次跳过渲染。最终交付前必须渲染；如确需跳过，在 accepted_risks 中加入 VAL-RENDER-SKIPPED。",
            )
        )
        manual_review.append("页面渲染尚未完成，不能确认视觉质量。")

    features = candidate_analysis["deck"]["package_features"]

    def add_manual_feature_check(check_id: str, name: str, details: str) -> None:
        checks.append(
            make_check(
                check_id,
                name,
                "pass" if check_id in accepted else "warning",
                (
                    f"{details} 已在 accepted_risks 中登记，仍需人工复核。"
                    if check_id in accepted
                    else f"{details} 未完成自动验证；请人工复核，并在 accepted_risks 中登记后再交付。"
                ),
            )
        )

    if features.get("diagram_parts"):
        add_manual_feature_check(
            "VAL-MANUAL-DIAGRAM",
            "SmartArt/Diagram 人工复核",
            "检测到 SmartArt/Diagram 部件。",
        )
        manual_review.append("检查 SmartArt/Diagram 的布局和可编辑性。")
    if features.get("embedded_objects"):
        add_manual_feature_check(
            "VAL-MANUAL-EMBEDDED",
            "嵌入对象人工复核",
            "检测到嵌入对象。",
        )
        manual_review.append("在 PowerPoint 中打开并检查所有嵌入对象。")
    if features.get("slide_timings"):
        add_manual_feature_check(
            "VAL-MANUAL-ANIMATION",
            "动画人工复核",
            "检测到动画或时间线。",
        )
        manual_review.append("播放检查动画顺序、触发和时序。")
    if features.get("external_relationships") or features.get("hyperlink_markers"):
        add_manual_feature_check(
            "VAL-MANUAL-LINKS",
            "链接人工复核",
            "检测到外部关系或超链接。",
        )
        manual_review.append("逐个测试外部链接和页内跳转。")

    statuses = [check["status"] for check in checks]
    unaccepted_warnings = [
        check for check in checks if check["status"] == "warning" and check["id"] not in accepted
    ]
    if "fail" in statuses or unaccepted_warnings:
        overall = "fail"
    elif "warning" in statuses or "not_checked" in statuses or manual_review:
        overall = "pass_with_accepted_risks"
    else:
        overall = "pass"

    return {
        "version": "1.0.0",
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "status": overall,
        "source": {
            "original": str(original_path),
            "original_sha256": sha256_file(original_path),
            "candidate": str(candidate_path),
            "candidate_sha256": sha256_file(candidate_path),
            "manifest": str(Path(manifest_path).resolve()),
        },
        "checks": checks,
        "differences": differences,
        "issues": issues,
        "manual_review_required": list(dict.fromkeys(manual_review)),
        "metrics": {
            "original_slide_count": original_count,
            "candidate_slide_count": candidate_count,
            "body_text_failure_slides": body_failures,
            "title_text_failure_slides": title_failures,
            "protected_slide_failures": protected_slide_failures,
            "numeric_failure_groups": numeric_failures,
            "protected_term_failures": protected_term_failures,
            "image_difference_slides": image_diffs,
            "candidate_p0_issues": len(p0_issues),
            "candidate_p1_issues": len(p1_issues),
            "render_status": render_status,
            "rendered_slide_count": rendered_count,
        },
    }


def markdown_report(report: dict[str, Any]) -> str:
    lines = [
        "# PPT 优化校验报告",
        "",
        f"- 结论：**{report['status']}**",
        f"- 原文件：`{report.get('source', {}).get('original', '')}`",
        f"- 候选文件：`{report.get('source', {}).get('candidate', '')}`",
        "",
        "## 校验项",
        "",
        "| ID | 校验项 | 状态 | 说明 |",
        "|---|---|---|---|",
    ]
    for check in report["checks"]:
        details = check["details"].replace("|", "\\|").replace("\n", " ")
        lines.append(f"| {check['id']} | {check['name']} | {check['status']} | {details} |")

    lines.extend(["", "## 差异", ""])
    if report["differences"]:
        for item in report["differences"]:
            lines.append(f"- `{item.get('type', 'difference')}`：{item}")
    else:
        lines.append("- 未检测到受保护内容差异。")

    lines.extend(["", "## 人工复核", ""])
    if report["manual_review_required"]:
        for item in report["manual_review_required"]:
            lines.append(f"- {item}")
    else:
        lines.append("- 无。")
    lines.append("")
    return "\n".join(lines)


def main() -> int:
    parser = argparse.ArgumentParser(description="Validate an optimized PPTX against the source and manifest.")
    parser.add_argument("--original", required=True)
    parser.add_argument("--candidate", required=True)
    parser.add_argument("--manifest", required=True)
    parser.add_argument("--out", required=True)
    parser.add_argument("--no-render", action="store_true")
    parser.add_argument("--visual-signoff", help="approved visual_signoff.json")
    parser.add_argument(
        "--allow-unconfirmed-visual",
        action="store_true",
        help="测试或结构审计专用：显式允许没有最终人工视觉确认；不可用于最终交付。",
    )
    args = parser.parse_args()

    try:
        out_dir = ensure_dir(args.out)
        report = validate_deck(
            args.original,
            args.candidate,
            args.manifest,
            out_dir,
            render=not args.no_render,
            visual_signoff_path=args.visual_signoff,
            allow_unconfirmed_visual=args.allow_unconfirmed_visual,
        )
        write_json(out_dir / "validation_report.json", report)
        (out_dir / "validation_report.md").write_text(markdown_report(report), encoding="utf-8")
        print(f"Validation status: {report['status']}")
        print(f"Wrote {out_dir / 'validation_report.json'}")
        return 0 if report["status"] != "fail" else 3
    except Exception as exc:  # noqa: BLE001
        print(f"ERROR: {exc}", file=sys.stderr)
        return 2


if __name__ == "__main__":
    raise SystemExit(main())
