#!/usr/bin/env python3
"""Create a small inconsistent deck for smoke testing the skill."""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def set_run(run, *, font: str, size: float, color: str, bold: bool = False) -> None:
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color.replace("#", ""))


def set_title(shape, text: str, *, font: str, size: float, color: str, x: float, y: float) -> None:
    shape.left = Inches(x)
    shape.top = Inches(y)
    shape.width = Inches(11.8)
    shape.height = Inches(0.65)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    set_run(run, font=font, size=size, color=color, bold=True)


def add_body(slide, lines: list[str], *, font: str, color: str = "000000") -> None:
    box = slide.shapes.add_textbox(Inches(0.9), Inches(1.55), Inches(11.6), Inches(4.8))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = line
        set_run(run, font=font, size=18 if idx == 0 else 16, color=color, bold=idx == 0)


def add_accent(slide, fill_color: str, line_color: str = "A5A5A5") -> None:
    shape = slide.shapes.add_shape(1, Inches(0.55), Inches(1.55), Inches(0.10), Inches(4.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill_color)
    shape.line.color.rgb = RGBColor.from_string(line_color)


def add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(10.5), Inches(7.08), Inches(2.1), Inches(0.22))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    set_run(run, font="Calibri", size=10, color="666666")


def build(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Cover
    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.background.fill.solid()
    cover.background.fill.fore_color.rgb = RGBColor.from_string("1F4E79")
    title = cover.shapes.title
    title.left = Inches(0.8)
    title.top = Inches(2.2)
    title.width = Inches(11.7)
    title.height = Inches(1.0)
    tf = title.text_frame
    tf.clear()
    run = tf.paragraphs[0].add_run()
    run.text = "智能工厂一体化平台方案"
    set_run(run, font="Arial", size=34, color="FFFFFF", bold=True)
    subtitle = cover.placeholders[1]
    subtitle.text = "PPT Review & Refinement Skill Smoke Test"
    for p in subtitle.text_frame.paragraphs:
        for r in p.runs:
            set_run(r, font="Calibri", size=18, color="FFFFFF")

    slides = [
        (
            "为什么需要统一平台",
            "Calibri",
            25,
            0.45,
            0.30,
            "5B9BD5",
            [
                "分散系统让管理动作和经营结果难以形成闭环",
                "需要同时连接 MES、WMS、QMS 和设备数据。",
                "核心指标包括 OTD、OEE 与一次合格率。",
            ],
        ),
        (
            "总体方案：业务、数据与AI协同",
            "Times New Roman",
            27,
            0.85,
            0.50,
            "70AD47",
            [
                "统一底座承载业务流程、数据对象和智能应用",
                "业务层覆盖计划、生产、质量、设备和仓储。",
                "数据层统一指标口径，AI 层面向具体工作流提供辅助决策。",
            ],
        ),
        (
            "下一步建议",
            "Arial",
            23,
            0.55,
            0.22,
            "5B9BD5",
            [
                "先完成蓝图与样板，再逐步扩展到全工厂",
                "第一阶段确认范围、主数据、接口和价值指标。",
                "第二阶段选择一个车间试点，并建立可复制模板。",
            ],
        ),
    ]

    for idx, (title_text, title_font, title_size, x, y, accent, body, ) in enumerate(slides, start=2):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        set_title(slide.shapes.title, title_text, font=title_font, size=title_size, color="1F4E79", x=x, y=y)
        add_accent(slide, accent)
        add_body(slide, body, font="Arial" if idx != 3 else "Calibri")
        add_footer(slide, f"{idx} / 4")

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--out", required=True)
    args = parser.parse_args()
    build(Path(args.out))
    print(f"Wrote {args.out}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
