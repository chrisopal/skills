#!/usr/bin/env python3
"""Regression tests for L2/L3 execution, report composition, and visual signoff."""

from __future__ import annotations

import json
import sys
import tempfile
import unittest
from pathlib import Path

from jsonschema import Draft202012Validator
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "scripts"))

from compose_review_report import compose_review_report  # noqa: E402
from confirm_visual_review import validate_visual_signoff  # noqa: E402
from execute_refinement_plan import execute_refinement_plan  # noqa: E402
from validate_pptx import validate_deck  # noqa: E402
from common import sha256_file  # noqa: E402


ROOT = Path(__file__).resolve().parents[1]


def build_deck(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    title.name = "Title Box"
    run = title.text_frame.paragraphs[0].add_run()
    run.text = "执行器测试标题"
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(0, 0, 0)
    prs.save(path)


def approved_l2_manifest() -> dict:
    manifest = json.loads((ROOT / "examples/change_manifest.l1.json").read_text(encoding="utf-8"))
    manifest["optimization_level"] = "L2"
    manifest["visual_policy"]["layout"] = "optimize"
    manifest["visual_policy"]["typography"] = "normalize"
    return manifest


def approved_pilot(source: Path, slide_count: int = 1) -> dict:
    return {
        "version": "1.0.0",
        "status": "approved",
        "reviewer": "Pilot Reviewer",
        "reviewed_at": "2026-08-07T00:00:00Z",
        "source_file": str(source),
        "source_sha256": sha256_file(source),
        "reviewed_slides": list(range(1, slide_count + 1)),
        "confirmed_tokens": ["typography.roles.page_title", "layout.title_box_in"],
        "notes": "样板已确认。",
    }


def issue(slide: int | str = 1) -> dict:
    return {
        "id": "NAR-001",
        "slide": slide,
        "scope": "slide" if isinstance(slide, int) else "deck",
        "category": "narrative",
        "severity": "P1",
        "finding": "主张没有在页首明确出现。",
        "impact": "读者需要额外寻找页面结论。",
        "recommendation": "将结论前置到标题或副标题。",
        "action_type": "rewrite_title",
        "approval_required": True,
        "confidence": "high",
        "evidence": ["title text"],
    }


class AdvancedCapabilityTests(unittest.TestCase):
    def test_report_composer_outputs_schema_valid_report(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            narrative = {
                "current_story": ["现状说明"],
                "proposed_story": ["结论先行", "证据支撑"],
                "global_issues": [issue("global")],
                "slides": [
                    {
                        "slide": 1,
                        "role": "cover",
                        "core_message": "建立统一主题",
                        "issues": [issue()],
                        "recommended_action": "强化主标题。",
                    }
                ],
            }
            visual = {
                "global_issues": [],
                "slides": [
                    {
                        "slide": 1,
                        "role": "cover",
                        "core_message": "建立统一主题",
                        "issues": [],
                        "recommended_action": "保持留白。",
                    }
                ],
            }
            scores = {
                key: 3
                for key in [
                    "narrative",
                    "theme",
                    "color",
                    "typography",
                    "hierarchy",
                    "layout",
                    "spacing",
                    "graphics",
                    "imagery",
                ]
            }
            for name, data in (("narrative.json", narrative), ("visual.json", visual), ("scores.json", scores)):
                (root / name).write_text(json.dumps(data, ensure_ascii=False), encoding="utf-8")
            output = root / "review_report.json"

            report = compose_review_report(
                source_file="source.pptx",
                narrative_path=root / "narrative.json",
                visual_path=root / "visual.json",
                scores_path=root / "scores.json",
                output_path=output,
                executive_summary="叙事和视觉方向已统一。",
                recommended_level="L2",
                manual_review_required=["逐页确认图表标签。"],
            )

            schema = json.loads((ROOT / "schemas/review_report.schema.json").read_text(encoding="utf-8"))
            self.assertEqual(report["recommended_level"], "L2")
            self.assertFalse(list(Draft202012Validator(schema).iter_errors(report)))

    def test_l2_executor_requires_approved_pilot_and_logs_editable_change(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "source.pptx"
            build_deck(source)
            manifest = root / "manifest.json"
            manifest.write_text(json.dumps(approved_l2_manifest()), encoding="utf-8")
            pilot = root / "pilot.json"
            pilot.write_text(json.dumps(approved_pilot(source)), encoding="utf-8")
            plan = root / "plan.json"
            plan.write_text(
                json.dumps(
                    {
                        "version": "1.0.0",
                        "source_file": "source.pptx",
                        "optimization_level": "L2",
                        "style_tokens_ref": str(ROOT / "examples/style_tokens.industrial-consulting.json"),
                        "change_manifest_ref": "manifest.json",
                        "slides": [
                            {
                                "slide": 1,
                                "archetype_current": "cover",
                                "archetype_target": "cover",
                                "objective": "统一标题位置和字号",
                                "keep": ["标题"],
                                "change": [
                                    {
                                        "target": "title",
                                        "action": "set_geometry",
                                        "parameters": {"shape_id": 2, "x": 2, "y": 1, "w": 7, "h": 1},
                                        "reason": "样板确认的标题框位置",
                                        "risk": "medium",
                                        "authorization": "pilot_confirmation",
                                    },
                                    {
                                        "target": "title",
                                        "action": "set_font_role",
                                        "parameters": {"shape_id": 2, "role": "page_title"},
                                        "reason": "采用确认的标题角色",
                                        "risk": "medium",
                                        "authorization": "pilot_confirmation",
                                    },
                                ],
                                "content_changes": [],
                                "image_actions": [],
                                "editability_requirements": ["标题保持可编辑"],
                                "manual_review_required": [],
                                "acceptance_checks": ["标题框位置正确"],
                            }
                        ],
                    },
                    ensure_ascii=False,
                ),
                encoding="utf-8",
            )
            output = root / "candidate.pptx"
            log_path = root / "execution_log.json"

            log = execute_refinement_plan(
                source,
                output,
                plan,
                ROOT / "examples/style_tokens.industrial-consulting.json",
                manifest,
                pilot,
                log_path=log_path,
                render=False,
            )

            self.assertGreaterEqual(log["change_count"], 2)
            self.assertTrue(output.exists())
            self.assertTrue(log_path.exists())
            reopened = Presentation(output)
            shape = next(shape for shape in reopened.slides[0].shapes if shape.shape_id == 2)
            self.assertAlmostEqual(shape.left / 914400, 2, places=3)
            self.assertEqual(shape.text_frame.paragraphs[0].runs[0].font.size.pt, 24)

            with self.assertRaises(PermissionError):
                execute_refinement_plan(
                    source,
                    root / "rejected.pptx",
                    plan,
                    ROOT / "examples/style_tokens.industrial-consulting.json",
                    manifest,
                    None,
                    log_path=root / "rejected.json",
                    render=False,
                )

            unsupported = json.loads(plan.read_text(encoding="utf-8"))
            unsupported["slides"][0]["change"][0]["action"] = "rebuild_chart"
            unsupported_plan = root / "unsupported-plan.json"
            unsupported_plan.write_text(json.dumps(unsupported), encoding="utf-8")
            with self.assertRaises(ValueError):
                execute_refinement_plan(
                    source,
                    root / "unsupported.pptx",
                    unsupported_plan,
                    ROOT / "examples/style_tokens.industrial-consulting.json",
                    manifest,
                    pilot,
                    log_path=root / "unsupported.json",
                    render=False,
                )

    def test_visual_signoff_is_required_by_final_validation(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            deck = root / "deck.pptx"
            build_deck(deck)
            manifest = root / "manifest.json"
            data = json.loads((ROOT / "examples/change_manifest.l1.json").read_text(encoding="utf-8"))
            data["accepted_risks"] = [{"id": "VAL-RENDER-SKIPPED", "description": "结构测试跳过渲染。"}]
            manifest.write_text(json.dumps(data), encoding="utf-8")
            report = validate_deck(deck, deck, manifest, root / "validation", render=False)
            self.assertEqual(report["status"], "fail")
            self.assertTrue(any(item["id"] == "VAL-VISUAL-SIGNOFF" for item in report["checks"]))

            signoff = root / "signoff.json"
            signoff.write_text(
                json.dumps(
                    {
                        "version": "1.0.0",
                        "status": "approved",
                        "reviewer": "Final Reviewer",
                        "reviewed_at": "2026-08-07T00:00:00Z",
                        "source_file": str(deck),
                        "source_sha256": sha256_file(deck),
                        "candidate_file": str(deck),
                        "candidate_sha256": sha256_file(deck),
                        "reviewed_slides": [1],
                        "confirmed_checks": [
                            "rendered_all_slides",
                            "checked_slide_overflow",
                            "checked_alignment",
                            "checked_fonts",
                            "checked_images",
                            "checked_editability",
                        ],
                        "notes": "最终人工视觉确认。",
                    },
                    ensure_ascii=False,
                ),
                encoding="utf-8",
            )
            self.assertTrue(validate_visual_signoff(signoff, deck, deck)["approved"])


if __name__ == "__main__":
    unittest.main()
