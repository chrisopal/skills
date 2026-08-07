#!/usr/bin/env python3
"""Regression tests for the PPT review/refinement quality gates."""

from __future__ import annotations

import json
import sys
import tempfile
import unittest
from pathlib import Path

from jsonschema import Draft202012Validator
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "scripts"))

from analyze_pptx import analyze_deck  # noqa: E402
from validate_pptx import validate_deck  # noqa: E402


ROOT = Path(__file__).resolve().parents[1]


def build_deck(path: Path, *, hyperlink: bool = False, overflowing_text: bool = False) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    width = Inches(1.0) if overflowing_text else Inches(8.0)
    height = Inches(0.2) if overflowing_text else Inches(1.0)
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), width, height)
    run = box.text_frame.paragraphs[0].add_run()
    run.text = (
        "这是一段很长的文字，放在极小的文本框中，按正常字号几乎必然会溢出页面布局。"
        if overflowing_text
        else "可验证页面"
    )
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(0, 0, 0)
    if hyperlink:
        run.hyperlink.address = "https://example.com"
    prs.save(path)


def example_manifest() -> dict:
    return json.loads((ROOT / "examples/change_manifest.l1.json").read_text(encoding="utf-8"))


class QualityGateTests(unittest.TestCase):
    def test_approved_manifest_requires_approver_identity(self) -> None:
        manifest = example_manifest()
        manifest["approval_status"] = "approved"
        manifest.pop("approved_by", None)
        manifest.pop("approved_at", None)
        schema = json.loads(
            (ROOT / "schemas/change_manifest.schema.json").read_text(encoding="utf-8")
        )

        errors = list(Draft202012Validator(schema).iter_errors(manifest))

        self.assertTrue(errors)

    def test_validation_rejects_invalid_manifest(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            deck = tmp_path / "source.pptx"
            manifest = tmp_path / "manifest.json"
            build_deck(deck)
            manifest.write_text(json.dumps({"version": "1.0.0"}), encoding="utf-8")

            with self.assertRaises(ValueError):
                validate_deck(deck, deck, manifest, tmp_path / "validation", render=False)

    def test_unaccepted_hyperlink_review_fails_validation(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            deck = tmp_path / "source.pptx"
            manifest = example_manifest()
            manifest["accepted_risks"] = [
                {
                    "id": "VAL-RENDER-SKIPPED",
                    "description": "测试中跳过渲染。",
                }
            ]
            manifest_path = tmp_path / "manifest.json"
            build_deck(deck, hyperlink=True)
            manifest_path.write_text(json.dumps(manifest), encoding="utf-8")

            report = validate_deck(deck, deck, manifest_path, tmp_path / "validation", render=False)

            self.assertEqual(report["status"], "fail")
            self.assertTrue(any(item["id"] == "VAL-MANUAL-LINKS" for item in report["checks"]))

    def test_accepted_hyperlink_review_is_reported_as_accepted_risk(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            deck = tmp_path / "source.pptx"
            manifest = example_manifest()
            manifest["accepted_risks"] = [
                {
                    "id": "VAL-RENDER-SKIPPED",
                    "description": "测试中跳过渲染。",
                },
                {
                    "id": "VAL-MANUAL-LINKS",
                    "description": "链接由人工测试。",
                },
            ]
            manifest_path = tmp_path / "manifest.json"
            build_deck(deck, hyperlink=True)
            manifest_path.write_text(json.dumps(manifest), encoding="utf-8")

            report = validate_deck(deck, deck, manifest_path, tmp_path / "validation", render=False)

            self.assertEqual(report["status"], "pass_with_accepted_risks")

    def test_unchanged_images_do_not_require_image_risk(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            deck = tmp_path / "source.pptx"
            manifest = example_manifest()
            manifest["image_policy"]["replace"] = "allowed"
            manifest["accepted_risks"] = [
                {
                    "id": "VAL-RENDER-SKIPPED",
                    "description": "测试中跳过渲染。",
                }
            ]
            manifest_path = tmp_path / "manifest.json"
            build_deck(deck)
            manifest_path.write_text(json.dumps(manifest), encoding="utf-8")

            report = validate_deck(deck, deck, manifest_path, tmp_path / "validation", render=False)

            image_check = next(check for check in report["checks"] if check["id"] == "VAL-IMAGES")
            self.assertEqual(image_check["status"], "pass")

    def test_analyzer_flags_obvious_text_overflow(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            deck = Path(tmp) / "overflow.pptx"
            build_deck(deck, overflowing_text=True)

            report = analyze_deck(deck)

            self.assertTrue(any(issue["id"].startswith("TXT-OVERFLOW") for issue in report["issues"]))


if __name__ == "__main__":
    unittest.main()
